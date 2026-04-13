#!/usr/bin/env python3
"""
Generate HIV Grand Rounds PPTX: MXM LA-ART & LA-PrEP Program
Task 49 - UCSF/ZSFG template-based presentation (23 slides)

Usage:
    python generate_deck.py [--output FILE]

Requires: python-pptx >= 1.0.0
"""

import argparse
import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ---------------------------------------------------------------------------
# Template layout indices (from UCSF_ZSFG_Template_16x9.pptx inspection)
# ---------------------------------------------------------------------------
LAYOUT_COVER_TEAL = 1      # Cover -- Teal (title + subtitle + date)
LAYOUT_BULLET = 3           # Bullet Slide (title + subtitle + content)
LAYOUT_BLANK = 4            # Blank (footer + slide number only)
LAYOUT_TWO_COLUMN = 6       # Two Column Slide (title + subtitle + 2 content)
LAYOUT_DIVIDER_NAVY = 10    # Divider Slide -- Navy
LAYOUT_DIVIDER_TEAL = 11    # Divider Slide -- Teal
LAYOUT_DIVIDER_BLUE = 12    # Divider Slide -- Blue
LAYOUT_BLANK_CLEAN = 13     # Blank (no placeholders)

# ---------------------------------------------------------------------------
# Color constants (UCSF brand)
# ---------------------------------------------------------------------------
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
TEAL = RGBColor(0x00, 0x95, 0xA8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
MID_GRAY = RGBColor(0x5F, 0x6B, 0x7A)
DARK_TEXT = RGBColor(0x1A, 0x20, 0x2C)
GOLD = RGBColor(0xD4, 0xA0, 0x1E)
LIGHT_TEAL = RGBColor(0xE0, 0xF7, 0xF9)
GREEN = RGBColor(0x32, 0xA0, 0x3E)
ORANGE = RGBColor(0xE8, 0x8D, 0x2A)
PURPLE = RGBColor(0x6B, 0x21, 0xA8)
RED_ACCENT = RGBColor(0xDC, 0x26, 0x26)
POLL_BG = RGBColor(0x05, 0x20, 0x49)  # Deep navy for poll slides


# ---------------------------------------------------------------------------
# Font constants
# ---------------------------------------------------------------------------
FONT_HEADING = "Garamond"
FONT_BODY = "Arial"
FONT_CODE = "Courier New"
HEADING_SIZE = Pt(36)
SUBHEADING_SIZE = Pt(18)
BODY_SIZE = Pt(22)
SMALL_SIZE = Pt(16)
CAPTION_SIZE = Pt(14)
STAT_SIZE = Pt(54)


# ---------------------------------------------------------------------------
# Helper functions
# ---------------------------------------------------------------------------
def set_speaker_notes(slide, notes_text):
    """Add speaker notes to a slide."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text


def set_bg_color(slide, color):
    """Set solid background color on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_name=FONT_BODY,
                font_size=BODY_SIZE, font_color=DARK_TEXT, bold=False,
                italic=False, alignment=PP_ALIGN.LEFT, word_wrap=True):
    """Add a styled text box and return the shape."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.italic = italic
    return txBox


def add_bullet_text(slide, left, top, width, height, bullets, font_name=FONT_BODY,
                    font_size=BODY_SIZE, font_color=DARK_TEXT, line_spacing=1.3,
                    bold_prefix=True):
    """Add a text box with multiple bullet lines. Supports **bold** prefix."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        p.line_spacing = line_spacing

        # Parse bold prefix: **text** rest
        if bold_prefix and bullet.startswith("**") and "**" in bullet[2:]:
            end_bold = bullet.index("**", 2)
            bold_text = bullet[2:end_bold]
            rest_text = bullet[end_bold+2:]

            run_bullet = p.add_run()
            run_bullet.text = "\u2022  "
            run_bullet.font.name = font_name
            run_bullet.font.size = font_size
            run_bullet.font.color.rgb = font_color

            run_bold = p.add_run()
            run_bold.text = bold_text
            run_bold.font.name = font_name
            run_bold.font.size = font_size
            run_bold.font.color.rgb = font_color
            run_bold.font.bold = True

            if rest_text.strip():
                run_rest = p.add_run()
                run_rest.text = rest_text
                run_rest.font.name = font_name
                run_rest.font.size = font_size
                run_rest.font.color.rgb = font_color
        else:
            run = p.add_run()
            run.text = f"\u2022  {bullet}"
            run.font.name = font_name
            run.font.size = font_size
            run.font.color.rgb = font_color

    return txBox


def add_stat_callout(slide, left, top, width, height, stat_text, label_text,
                     stat_color=TEAL, label_color=DARK_TEXT):
    """Add a large stat number with label below."""
    # Stat number
    add_textbox(slide, left, top, width, height * 0.6, stat_text,
                font_name=FONT_HEADING, font_size=STAT_SIZE,
                font_color=stat_color, bold=True, alignment=PP_ALIGN.CENTER)
    # Label
    add_textbox(slide, left, top + height * 0.55, width, height * 0.4, label_text,
                font_name=FONT_BODY, font_size=SMALL_SIZE,
                font_color=label_color, alignment=PP_ALIGN.CENTER)


def add_divider_slide(prs, title, subtitle=""):
    """Add a teal divider slide using the template layout."""
    layout = prs.slide_layouts[LAYOUT_DIVIDER_TEAL]
    slide = prs.slides.add_slide(layout)
    # Set title placeholder
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = title
            for para in ph.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = FONT_HEADING
                    run.font.size = HEADING_SIZE
    if subtitle:
        add_textbox(slide, 1.0, 4.8, 11.0, 0.6, subtitle,
                    font_size=SUBHEADING_SIZE, font_color=WHITE,
                    alignment=PP_ALIGN.CENTER)
    return slide


def add_navy_content_slide(prs, title):
    """Add a blank slide with navy background and white title."""
    layout = prs.slide_layouts[LAYOUT_BLANK_CLEAN]
    slide = prs.slides.add_slide(layout)
    set_bg_color(slide, NAVY)
    add_textbox(slide, 0.75, 0.3, 11.8, 0.8, title,
                font_name=FONT_HEADING, font_size=HEADING_SIZE,
                font_color=WHITE, bold=True)
    return slide


def add_white_content_slide(prs, title):
    """Add a blank slide with white background and navy title."""
    layout = prs.slide_layouts[LAYOUT_BLANK_CLEAN]
    slide = prs.slides.add_slide(layout)
    add_textbox(slide, 0.75, 0.3, 11.8, 0.8, title,
                font_name=FONT_HEADING, font_size=HEADING_SIZE,
                font_color=NAVY, bold=True)
    return slide


def add_poll_slide(prs, question, options, correct_idx=None):
    """Add a poll slide with navy background, question, and lettered options."""
    layout = prs.slide_layouts[LAYOUT_BLANK_CLEAN]
    slide = prs.slides.add_slide(layout)
    set_bg_color(slide, POLL_BG)

    # Poll header badge
    add_textbox(slide, 0.75, 0.3, 3.5, 0.5, "POLL EVERYWHERE",
                font_name=FONT_BODY, font_size=CAPTION_SIZE,
                font_color=TEAL, bold=True)

    # Question
    add_textbox(slide, 0.75, 0.9, 11.8, 1.0, question,
                font_name=FONT_HEADING, font_size=Pt(28),
                font_color=WHITE, bold=True)

    # Options
    letters = "ABCDE"
    y_start = 2.2
    for i, opt in enumerate(options):
        option_color = WHITE
        prefix = f"{letters[i]}:  "
        txBox = slide.shapes.add_textbox(
            Inches(1.2), Inches(y_start + i * 0.85),
            Inches(11.0), Inches(0.7)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]

        run_letter = p.add_run()
        run_letter.text = prefix
        run_letter.font.name = FONT_BODY
        run_letter.font.size = Pt(20)
        run_letter.font.color.rgb = TEAL
        run_letter.font.bold = True

        run_text = p.add_run()
        run_text.text = opt
        run_text.font.name = FONT_BODY
        run_text.font.size = Pt(20)
        run_text.font.color.rgb = option_color

    return slide


def add_two_column_navy(prs, title, left_title, left_bullets, right_title, right_bullets):
    """Navy background, two-column content with sub-headers."""
    slide = add_navy_content_slide(prs, title)

    # Left column sub-header
    add_textbox(slide, 0.75, 1.2, 5.8, 0.5, left_title,
                font_name=FONT_BODY, font_size=SUBHEADING_SIZE,
                font_color=TEAL, bold=True)
    # Left bullets
    add_bullet_text(slide, 0.75, 1.8, 5.8, 5.0, left_bullets,
                    font_size=SMALL_SIZE, font_color=WHITE, line_spacing=1.2)

    # Right column sub-header
    add_textbox(slide, 7.0, 1.2, 5.8, 0.5, right_title,
                font_name=FONT_BODY, font_size=SUBHEADING_SIZE,
                font_color=TEAL, bold=True)
    # Right bullets
    add_bullet_text(slide, 7.0, 1.8, 5.8, 5.0, right_bullets,
                    font_size=SMALL_SIZE, font_color=WHITE, line_spacing=1.2)

    return slide


def add_two_column_white(prs, title, left_title, left_bullets, right_title, right_bullets):
    """White background, two-column content with sub-headers."""
    slide = add_white_content_slide(prs, title)

    add_textbox(slide, 0.75, 1.2, 5.8, 0.5, left_title,
                font_name=FONT_BODY, font_size=SUBHEADING_SIZE,
                font_color=TEAL, bold=True)
    add_bullet_text(slide, 0.75, 1.8, 5.8, 5.0, left_bullets,
                    font_size=SMALL_SIZE, font_color=DARK_TEXT, line_spacing=1.2)

    add_textbox(slide, 7.0, 1.2, 5.8, 0.5, right_title,
                font_name=FONT_BODY, font_size=SUBHEADING_SIZE,
                font_color=TEAL, bold=True)
    add_bullet_text(slide, 7.0, 1.8, 5.8, 5.0, right_bullets,
                    font_size=SMALL_SIZE, font_color=DARK_TEXT, line_spacing=1.2)

    return slide


def add_table(slide, headers, rows, left=0.75, top=2.0, width=11.8,
              font_color=DARK_TEXT, header_bg=NAVY, header_fg=WHITE):
    """Add a formatted table to a slide."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    row_height = Inches(0.45)

    table_shape = slide.shapes.add_table(
        n_rows, n_cols, Inches(left), Inches(top),
        Inches(width), row_height * n_rows
    )
    table = table_shape.table

    # Header row
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_bg
        for para in cell.text_frame.paragraphs:
            para.alignment = PP_ALIGN.LEFT
            for run in para.runs:
                run.font.name = FONT_BODY
                run.font.size = Pt(14)
                run.font.bold = True
                run.font.color.rgb = header_fg

    # Data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, value in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = FONT_BODY
                    run.font.size = Pt(13)
                    run.font.color.rgb = font_color

    return table_shape


# ---------------------------------------------------------------------------
# Patient 2 Timeline Figure
# ---------------------------------------------------------------------------
def add_timeline_figure(slide):
    """Draw a horizontal timeline for Patient 2 with dose markers and annotations."""
    # Timeline arrow (horizontal line)
    arrow_y = 3.5
    arrow_left = 0.8
    arrow_width = 11.5

    # Background band
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(arrow_left), Inches(arrow_y - 0.05),
        Inches(arrow_width), Inches(0.1)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()

    # Timeline phases with dose markers
    # Phase 1: W-29 to W0 (loading + escalation)
    # Phase 2: W0 decision point
    # Phase 3: W38 return, W42 restart, W107 publication

    phases = [
        {"x": 1.0, "label": "W-29", "detail": "CAB/RPV\nLoading", "color": PURPLE, "below": "VL 177,992\nCD4 187"},
        {"x": 2.8, "label": "W-24", "detail": "Doses\n#2-5", "color": GREEN, "below": "VL 33\u219230"},
        {"x": 4.0, "label": "W-17", "detail": "Delayed\n12d", "color": ORANGE, "below": ""},
        {"x": 5.0, "label": "W-8", "detail": "Delayed\n21d", "color": ORANGE, "below": "VL <30"},
        {"x": 6.5, "label": "W0", "detail": "VIREMIA\nVL 37,464", "color": RED_ACCENT, "below": "Added LEN\n+ CAB/RPV"},
        {"x": 8.2, "label": "W38", "detail": "Returns\nabscess", "color": GREEN, "below": "VL <30\n(LEN mono!)"},
        {"x": 9.5, "label": "W42", "detail": "Restart\nLEN+CAB", "color": GREEN, "below": "VL UD"},
        {"x": 11.0, "label": "W107", "detail": "Published", "color": GREEN, "below": "CD4 554\n65+ wks sup"},
    ]

    for phase in phases:
        x = phase["x"]
        color = phase["color"]

        # Dose circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x - 0.12), Inches(arrow_y - 0.15),
            Inches(0.3), Inches(0.3)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()

        # Label above (week number)
        add_textbox(slide, x - 0.5, arrow_y - 0.9, 1.0, 0.35, phase["label"],
                    font_name=FONT_BODY, font_size=Pt(10), font_color=WHITE,
                    bold=True, alignment=PP_ALIGN.CENTER)

        # Detail above
        if phase["detail"]:
            add_textbox(slide, x - 0.7, arrow_y - 1.6, 1.4, 0.65, phase["detail"],
                        font_name=FONT_BODY, font_size=Pt(9), font_color=LIGHT_GRAY,
                        alignment=PP_ALIGN.CENTER)

        # Lab values below
        if phase["below"]:
            add_textbox(slide, x - 0.7, arrow_y + 0.3, 1.4, 0.65, phase["below"],
                        font_name=FONT_CODE, font_size=Pt(9), font_color=TEAL,
                        alignment=PP_ALIGN.CENTER)

    # Gold callout boxes for key decisions
    # Decision at W0
    callout = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.5), Inches(4.7), Inches(3.2), Inches(0.6)
    )
    callout.fill.solid()
    callout.fill.fore_color.rgb = GOLD
    callout.line.fill.background()
    tf = callout.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Pre-genotype: Added LEN as 2nd active agent"
    run.font.name = FONT_BODY
    run.font.size = Pt(10)
    run.font.color.rgb = NAVY
    run.font.bold = True

    # Legend
    legend_y = 6.2
    legend_items = [
        (PURPLE, "Initiation"), (GREEN, "On-time"), (ORANGE, "Delayed"), (RED_ACCENT, "Viremia")
    ]
    for i, (color, label) in enumerate(legend_items):
        x = 3.0 + i * 2.2
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x), Inches(legend_y), Inches(0.2), Inches(0.2)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()
        add_textbox(slide, x + 0.3, legend_y - 0.05, 1.5, 0.3, label,
                    font_name=FONT_BODY, font_size=Pt(11), font_color=WHITE)


# ---------------------------------------------------------------------------
# Build all 23 slides
# ---------------------------------------------------------------------------
def remove_template_slides(prs):
    """Remove all pre-existing demo slides from the template."""
    # python-pptx stores slide references in the XML sldIdLst
    xml_slides = prs.slides._sldIdLst
    slide_ids = list(xml_slides)
    for sldId in slide_ids:
        rId = sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        prs.part.drop_rel(rId)
        xml_slides.remove(sldId)


def build_deck(template_path, output_path):
    """Build the complete 23-slide HIV Grand Rounds presentation."""
    prs = Presentation(template_path)
    remove_template_slides(prs)

    # ======================================================================
    # SLIDE 1: Title / Section Opener (Divider Teal)
    # ======================================================================
    slide1 = add_divider_slide(prs, "LA-ART & LA-PrEP at MXM",
                                "HIV Grand Rounds, UCSF/ZSFG")
    add_textbox(slide1, 1.0, 5.5, 11.0, 0.4, "Nicky Mehtani, MD MPH  |  April 2026",
                font_size=SUBHEADING_SIZE, font_color=WHITE,
                alignment=PP_ALIGN.CENTER)
    set_speaker_notes(slide1,
        "Thanks Stephen for that overview. I'm going to walk through our LA-ART and "
        "LA-PrEP data, then we'll get into four patient cases. I've got some polls "
        "set up -- pull out your phones for those.")

    # ======================================================================
    # SLIDE 2: LA-ART Who We Serve (Two Column Navy)
    # ======================================================================
    slide2 = add_two_column_navy(prs,
        "LA-ART \u2014 Who We Serve",
        "Demographics (n=34, Nov 2021 \u2013 Apr 2026)",
        [
            "Median age: 40",
            "22 cis-men, 8 trans women/NB, 4 cis-women",
            "50% Black, 26% Hispanic, 24% non-Hispanic White",
            "10 street homeless, 16 sheltered/couch-surfing, 8 unstably housed",
        ],
        "Clinical Complexity at Initiation",
        [
            "85% methamphetamine use disorder",
            "24% opioid use disorder",
            "32% schizophrenia or bipolar with psychosis",
            "32% initiated with CD4 <200 (7 with CD4 <50)",
            "71% started with detectable viremia",
            "Of those, 17/24 NOT suppressed \u22653 years",
        ]
    )
    set_speaker_notes(slide2,
        "So this is who we're serving. Thirty-four patients since November 2021. "
        "The majority are people with active methamphetamine use, serious mental illness, "
        "and unstable housing. Most had detectable virus when we started -- and most of "
        "those had been viremic for years. These are not people who were already doing "
        "well on oral ART.")

    # ======================================================================
    # SLIDE 3: LA-ART Delivery & Outcomes (Navy with stat callouts)
    # ======================================================================
    slide3 = add_navy_content_slide(prs, "LA-ART \u2014 Delivery & Outcomes")

    # Stat callouts
    add_stat_callout(slide3, 0.5, 1.2, 4.0, 1.8, "100%", "virally suppressed\nat latest follow-up",
                     stat_color=TEAL, label_color=WHITE)
    add_stat_callout(slide3, 4.5, 1.2, 4.0, 1.8, "93%", "injections on time\n(within 7-day window)",
                     stat_color=GREEN, label_color=WHITE)
    add_stat_callout(slide3, 8.5, 1.2, 4.0, 1.8, "747", "total injection visits",
                     stat_color=GOLD, label_color=WHITE)

    # Details below stats
    details = [
        "80% at MXM clinic | 17% mobile outreach, shelters, syringe access | 3% jails/hospitals",
        "Median time on LA-ART: 116 weeks (IQR: 60\u2013145; range: 2\u2013209 weeks)",
        "9 successfully transferred to new primary care homes",
        "1 currently LTFU ~3 months (Patient 1)",
    ]
    add_bullet_text(slide3, 0.75, 3.5, 11.8, 2.5, details,
                    font_size=SMALL_SIZE, font_color=WHITE, line_spacing=1.3)

    # Current regimens
    add_textbox(slide3, 0.75, 5.5, 4.0, 0.4, "Current Regimens:",
                font_size=SMALL_SIZE, font_color=TEAL, bold=True)
    regimens = [
        "CAB/RPV q1-month: n=12",
        "CAB/RPV q2-month: n=11",
        "CAB (\u00b1RPV) q1-mo + LEN q6-mo: n=9",
        "CAB/RPV q2-mo + LEN q6-mo: n=2",
    ]
    add_bullet_text(slide3, 0.75, 5.9, 11.8, 1.5, regimens,
                    font_size=CAPTION_SIZE, font_color=LIGHT_GRAY, line_spacing=1.1)

    set_speaker_notes(slide3,
        "Seven hundred and forty-seven injection visits. The vast majority happen at "
        "the clinic, but we're doing about 17% in the field -- shelters, syringe access "
        "programs, mobile outreach. Ninety-three percent on time. And the headline: "
        "100% virally suppressed. Every single one of our 34 patients has an undetectable "
        "viral load at last follow-up. Median time on therapy is over two years. We've "
        "also successfully transferred 9 patients to new primary care homes -- this isn't "
        "meant to be forever for everyone.")

    # ======================================================================
    # SLIDE 4: Oral vs. LA-ART Comparison (Two Column White with table)
    # ======================================================================
    slide4 = add_white_content_slide(prs, "Oral vs. LA-ART Comparison")
    add_textbox(slide4, 0.75, 1.1, 6.0, 0.4, "2023\u20132024, Mehtani et al., in prep",
                font_size=CAPTION_SIZE, font_color=MID_GRAY, italic=True)

    add_table(slide4,
        ["Outcome", "Oral ART (n=123)", "LA-ART (n=24)"],
        [
            ["Ever virally suppressed", "89/123 (72%)", "24/24 (100%)"],
            ["Suppressed at last follow-up", "69/123 (56%)", "24/24 (100%)"],
            ["Mean encounters/person-year", "9.1", "23.0"],
        ],
        top=1.8, width=11.0
    )

    # Caveat box
    caveat_box = slide4.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.75), Inches(4.2), Inches(11.5), Inches(1.8)
    )
    caveat_box.fill.solid()
    caveat_box.fill.fore_color.rgb = RGBColor(0xFF, 0xF8, 0xE1)  # light gold
    caveat_box.line.color.rgb = GOLD
    caveat_box.line.width = Pt(2)
    tf = caveat_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run_label = p.add_run()
    run_label.text = "Important Caveat:  "
    run_label.font.name = FONT_BODY
    run_label.font.size = SMALL_SIZE
    run_label.font.color.rgb = NAVY
    run_label.font.bold = True
    run_text = p.add_run()
    run_text.text = ("Groups not randomized. LA-ART patients selected because we could "
                     "reliably locate them and they demonstrated interest/understanding of "
                     "injection requirements. Outcome gap reflects BOTH the intervention "
                     "AND the selection.")
    run_text.font.name = FONT_BODY
    run_text.font.size = SMALL_SIZE
    run_text.font.color.rgb = DARK_TEXT

    set_speaker_notes(slide4,
        "Here's a head-to-head comparison from our manuscript in preparation. Among 128 "
        "people with HIV at MXM, the LA-ART group had 100% suppression -- both ever and "
        "at last follow-up. Oral ART was 72% ever suppressed and only 56% at last check. "
        "But -- and this is critical -- these groups are not randomized. We selected LA-ART "
        "patients specifically because we could find them reliably and they understood the "
        "injection commitment. So the gap reflects both the intervention and the selection. "
        "This is not a randomized trial.")

    # ======================================================================
    # SLIDE 5: LA-PrEP at MXM (White content)
    # ======================================================================
    slide5 = add_white_content_slide(prs, "LA-PrEP at MXM")

    # Stat callouts top row
    add_stat_callout(slide5, 0.5, 1.1, 3.5, 1.5, "68", "prescribed",
                     stat_color=TEAL, label_color=DARK_TEXT)
    add_stat_callout(slide5, 3.5, 1.1, 3.5, 1.5, "52", "initiated (76%)",
                     stat_color=TEAL, label_color=DARK_TEXT)
    add_stat_callout(slide5, 7.0, 1.1, 3.0, 1.5, "90%", "retained \u22653 months",
                     stat_color=GREEN, label_color=DARK_TEXT)
    add_stat_callout(slide5, 10.0, 1.1, 3.0, 1.5, "2", "seroconversions",
                     stat_color=RED_ACCENT, label_color=DARK_TEXT)

    prep_bullets = [
        "Cohort: 65% unhoused, 55% stimulant use, 27% OUD, 19% TG/NB, 50% non-White",
        "304 injections: 58% clinic, 21% shelters, 23% mobile outreach",
        "Timing: 78% on time, 17% delayed 8\u201328 days, 7% delayed >28 days",
        "HIV testing: 68% any testing; 61% CDC guideline-concordant blood-based",
        "Both seroconversions during prolonged lapses, NOT during imperfect dosing",
        "Gift card disruption (~6 months) likely explains suboptimal retention",
    ]
    add_bullet_text(slide5, 0.75, 3.0, 11.8, 3.5, prep_bullets,
                    font_size=SMALL_SIZE, font_color=DARK_TEXT, line_spacing=1.25)

    # Key message box
    msg_box = slide5.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.75), Inches(5.8), Inches(11.5), Inches(0.9)
    )
    msg_box.fill.solid()
    msg_box.fill.fore_color.rgb = LIGHT_TEAL
    msg_box.line.color.rgb = TEAL
    msg_box.line.width = Pt(2)
    tf = msg_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = ("Continuity of coverage > protocol perfection. For PEH/PWUD: "
                "prioritize rapid initiation, outreach-based delivery, flexible lab protocols.")
    run.font.name = FONT_BODY
    run.font.size = SMALL_SIZE
    run.font.color.rgb = NAVY
    run.font.bold = True

    set_speaker_notes(slide5,
        "LA-PrEP -- sixty-eight prescribed, fifty-two initiated. Ninety percent retained "
        "at three months, which is actually remarkable for this population. Two seroconversions, "
        "and both happened during prolonged lapses in coverage -- not during imperfect dosing. "
        "That's the key message here: continuity of coverage matters more than protocol perfection. "
        "We lost our gift card incentives for about 6 months -- that almost certainly explains "
        "some of the retention drop. We're planning to write that up. One additional seroconversion "
        "happened after a patient transferred to TWUHC -- not ours to present formally, but worth "
        "knowing about.")

    # ======================================================================
    # SLIDE 6: Case Divider - Patient 1
    # ======================================================================
    slide6 = add_divider_slide(prs, "Patient 1",
                                "49M, SRO housing, schizophrenia, methamphetamine use disorder")
    set_speaker_notes(slide6,
        "Alright, cases. Four patients. Get your phones out for Poll Everywhere. "
        "Patient 1 is published as Patient B in our Open Forum Infectious Diseases paper.")

    # ======================================================================
    # SLIDE 7: Patient 1 Background
    # ======================================================================
    slide7 = add_white_content_slide(prs, "Patient 1 \u2014 Background")
    p1_bullets = [
        "49-year-old man, housed in SRO (Tenderloin)",
        "HIV diagnosed 2010",
        "Schizophrenia, bipolar disorder, psychosis NOS; methamphetamine use disorder",
        "NOT interested in antipsychotics or psychiatry; no cell phone",
        "Unable to maintain oral ART x5 years (tried EFV/TDF/FTC, RPV/TDF/FTC, EVG/c/TAF/FTC, BIC/TAF/FTC)",
        "Baseline: CD4 334, VL 340 c/mL",
        "Genotype: Y181C + M184V = high-level RPV resistance; NO INSTI mutations",
        "3rd visit to MXM in 3 months \u2014 first sustained engagement in years",
    ]
    add_bullet_text(slide7, 0.75, 1.3, 11.8, 5.5, p1_bullets,
                    font_size=Pt(18), font_color=DARK_TEXT, line_spacing=1.25)
    set_speaker_notes(slide7,
        "Patient 1. Forty-nine years old, living in an SRO in the Tenderloin. HIV since 2010. "
        "Schizophrenia, bipolar, psychosis -- not interested in psychiatric medications. No cell "
        "phone. He'd been through four oral regimens, couldn't stick with any of them. When he "
        "walked into MXM it was his third visit in three months -- the first time he'd been "
        "consistently showing up anywhere in years. The key thing to know about his resistance: "
        "Y181C gives him high-level RPV resistance. M184V. But no INSTI mutations -- so CAB is still active.")

    # ======================================================================
    # SLIDE 8: Patient 1 Poll 1
    # ======================================================================
    slide8 = add_poll_slide(prs,
        "Would you start this patient on LA-ART?",
        [
            "No \u2014 not with untreated schizophrenia/psychosis",
            "Yes, but wait until he has a cell phone",
            "Yes \u2014 start CAB/RPV monthly",
            "Yes \u2014 start CAB monthly + LEN q6 months",
            "Yes \u2014 start CAB/RPV q2-months + LEN q6 months",
        ],
        correct_idx=3
    )
    set_speaker_notes(slide8,
        "Pull up Poll Everywhere. Would you start this patient on LA-ART? And if so, what regimen?")

    # ======================================================================
    # SLIDE 9: Patient 1 Poll 1 Answer + Clinical Course
    # ======================================================================
    slide9 = add_white_content_slide(prs, "Patient 1 \u2014 Why D Is Correct")

    answers = [
        "A: Untreated psychiatric illness is the rule, not exception, at MXM",
        "B: No phone is not a barrier \u2014 outreach RN goes to his SRO",
        "C: Y181C = high-level RPV resistance; CAB monotherapy risk with any lapse",
        "E: Including RPV with documented resistance weakens regimen",
    ]
    add_bullet_text(slide9, 0.75, 1.2, 11.8, 2.0, answers,
                    font_size=CAPTION_SIZE, font_color=MID_GRAY, line_spacing=1.2)

    # Correct answer callout
    d_box = slide9.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.75), Inches(3.0), Inches(11.5), Inches(0.7)
    )
    d_box.fill.solid()
    d_box.fill.fore_color.rgb = LIGHT_TEAL
    d_box.line.color.rgb = TEAL
    d_box.line.width = Pt(2)
    tf = d_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "D: Drop RPV (resistance), add LEN as capsid inhibitor backbone to protect INSTI class during lapses"
    run.font.name = FONT_BODY
    run.font.size = Pt(16)
    run.font.color.rgb = NAVY
    run.font.bold = True

    # Phase 1 success
    add_textbox(slide9, 0.75, 4.0, 6.0, 0.4, "Phase 1 \u2014 Success (14 months):",
                font_size=SUBHEADING_SIZE, font_color=TEAL, bold=True)
    phase1_bullets = [
        "CAB monthly + LEN q6 months",
        "16 injections, ALL on time (MXM or outreach RN at SRO)",
        "VL undetectable throughout",
    ]
    add_bullet_text(slide9, 0.75, 4.5, 5.5, 1.5, phase1_bullets,
                    font_size=CAPTION_SIZE, font_color=DARK_TEXT, line_spacing=1.2)

    # Then: LTFU
    add_textbox(slide9, 7.0, 4.0, 5.5, 0.4, "Then: LTFU",
                font_size=SUBHEADING_SIZE, font_color=RED_ACCENT, bold=True)
    ltfu_bullets = [
        "Disappeared ~3 months; biked >300 miles",
        "Returns 10 wks after last CAB, 4 wks after last LEN",
        "Wants to restart",
    ]
    add_bullet_text(slide9, 7.0, 4.5, 5.5, 1.5, ltfu_bullets,
                    font_size=CAPTION_SIZE, font_color=DARK_TEXT, line_spacing=1.2)

    set_speaker_notes(slide9,
        "The answer is D. Drop RPV -- he has high-level resistance. Start CAB monthly and add "
        "lenacapavir every six months. LEN gives you a second active agent from a completely "
        "different drug class. If he lapses -- and in this population, lapses happen -- LEN protects "
        "the INSTI class while CAB levels drop. And that's exactly what happened. Fourteen months of "
        "success -- every injection on time, viral load undetectable. Then he disappeared. Turns out "
        "he biked over 300 miles to another city. Came back 10 weeks after his last CAB was due. "
        "Four weeks after his last LEN was due. And he wanted to restart.")

    # ======================================================================
    # SLIDE 10: Patient 1 Poll 2
    # ======================================================================
    slide10 = add_poll_slide(prs,
        "He returns 10 weeks after last CAB/RPV, 4 weeks after last LEN.\n"
        "In addition to drawing HIV VL + RNA genotype, what do you do today?",
        [
            "Nothing today \u2014 have him return in 2\u20133 days once labs result",
            "Start BIC/TAF/FTC oral bridge until labs return",
            "Start DRV/c/TAF/FTC empirically until labs return",
            "Administer LEN + CAB today and plan for close follow-up",
            "Another reasonable option? (open to audience)",
        ]
    )
    set_speaker_notes(slide10,
        "Poll 2. He's standing in front of you right now. You've drawn labs. What do you do "
        "TODAY? No single correct answer here -- this is a genuine discussion question.")

    # ======================================================================
    # SLIDE 11: Patient 1 Resolution + LEN Pharmacology
    # ======================================================================
    slide11 = add_white_content_slide(prs, "Patient 1 \u2014 Resolution")

    # Labs returned
    labs_box = slide11.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.75), Inches(1.2), Inches(11.5), Inches(0.6)
    )
    labs_box.fill.solid()
    labs_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)  # light green
    labs_box.line.color.rgb = GREEN
    labs_box.line.width = Pt(2)
    tf = labs_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Labs returned: VL UNDETECTABLE despite 10 weeks late for CAB, 4 weeks late for LEN"
    run.font.name = FONT_BODY
    run.font.size = Pt(18)
    run.font.color.rgb = NAVY
    run.font.bold = True

    # Teaching points
    add_textbox(slide11, 0.75, 2.1, 5.5, 0.4, "Teaching Point 1: Low-barrier philosophy",
                font_size=SUBHEADING_SIZE, font_color=TEAL, bold=True)
    add_bullet_text(slide11, 0.75, 2.5, 5.5, 1.2,
        ["Act with the opportunity in front of you",
         "Delay = further delay = months of lost contact"],
        font_size=CAPTION_SIZE, font_color=DARK_TEXT, line_spacing=1.2)

    add_textbox(slide11, 7.0, 2.1, 5.5, 0.4, "Teaching Point 2: Why LEN protected him",
                font_size=SUBHEADING_SIZE, font_color=TEAL, bold=True)
    add_bullet_text(slide11, 7.0, 2.5, 5.5, 1.2,
        ["LEN has longer PK tail than CAB",
         "Capsid inhibitor = different mechanism",
         "Without LEN, INSTI resistance almost certain"],
        font_size=CAPTION_SIZE, font_color=DARK_TEXT, line_spacing=1.2)

    # Phase 2 success
    add_textbox(slide11, 0.75, 4.0, 6.0, 0.4, "Phase 2 \u2014 Success (again):",
                font_size=SUBHEADING_SIZE, font_color=GREEN, bold=True)
    add_bullet_text(slide11, 0.75, 4.4, 11.8, 1.2,
        ["LEN + CAB monthly restarted; 16 more injections over ~16 months, all on time",
         "VL undetectable throughout"],
        font_size=CAPTION_SIZE, font_color=DARK_TEXT, line_spacing=1.2)

    # Publication
    add_textbox(slide11, 0.75, 5.4, 11.8, 0.4,
                "Published: Mehtani NJ et al., Open Forum Infect Dis 2025;12(6):ofaf330. PMC12188974.",
                font_size=Pt(12), font_color=MID_GRAY, italic=True)

    set_speaker_notes(slide11,
        "His viral load came back undetectable. Ten weeks late for CAB and he was still "
        "suppressed. Two teaching points. First -- the low-barrier philosophy. He's a guy "
        "who bikes 300 miles on a whim. If we'd sent him away to wait for labs, we might "
        "not see him for another three months. You act when the patient is in front of you. "
        "Second -- why was he still suppressed? LEN has a longer PK tail than CAB. When CAB "
        "dropped to subtherapeutic levels, LEN was likely still working. Different mechanism "
        "-- capsid inhibitor versus integrase inhibitor. Without LEN in this regimen, I'm "
        "almost certain he would have developed INSTI resistance during that 10-week gap.")

    # ======================================================================
    # SLIDE 12: Patient 1 LTFU + Population Question
    # ======================================================================
    slide12 = add_navy_content_slide(prs, "Patient 1 \u2014 Current Status & Population Question")

    add_textbox(slide12, 0.75, 1.2, 5.5, 0.4, "Current Status:",
                font_size=SUBHEADING_SIZE, font_color=TEAL, bold=True)
    status_bullets = [
        "Disappeared again after 16 more months of success",
        "11 weeks since last CAB/RPV dose",
        "Still within LEN window (due ~26 weeks)",
        "Cannot locate \u2014 suspected >3,000 miles away",
        "**Only patient in our cohort currently LTFU**",
    ]
    add_bullet_text(slide12, 0.75, 1.7, 5.5, 3.0, status_bullets,
                    font_size=SMALL_SIZE, font_color=WHITE, line_spacing=1.2)

    # Population question box
    q_box = slide12.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.8), Inches(1.2), Inches(5.8), Inches(3.5)
    )
    q_box.fill.solid()
    q_box.fill.fore_color.rgb = RGBColor(0x0D, 0x30, 0x5C)  # slightly lighter navy
    q_box.line.color.rgb = TEAL
    q_box.line.width = Pt(2)
    tf = q_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Population-Level Question"
    run.font.name = FONT_HEADING
    run.font.size = Pt(20)
    run.font.color.rgb = TEAL
    run.font.bold = True
    p2 = tf.add_paragraph()
    p2.space_before = Pt(12)
    run2 = p2.add_run()
    run2.text = ("Should LEN + CAB/RPV become the DEFAULT for high LTFU-risk patients "
                 "\u2014 as harm reduction to protect the INSTI class?")
    run2.font.name = FONT_BODY
    run2.font.size = SMALL_SIZE
    run2.font.color.rgb = WHITE
    p3 = tf.add_paragraph()
    p3.space_before = Pt(12)
    run3 = p3.add_run()
    run3.text = ("\u22652 patients at other counties started on CAB/RPV alone, later presented "
                 "with new INSTI resistance after LTFU.")
    run3.font.name = FONT_BODY
    run3.font.size = CAPTION_SIZE
    run3.font.color.rgb = GOLD
    run3.font.italic = True

    set_speaker_notes(slide12,
        "He's disappeared again. Eleven weeks since his last CAB. But he's still within the "
        "LEN window. We think he's over 3,000 miles away based on a Care Everywhere ED note. "
        "He's the only patient in our cohort who's LTFU. The silver lining -- if he develops "
        "resistance, it's more likely to be LEN resistance than INSTI resistance. And LEN "
        "resistance mutations confer very low viral fitness. The bigger question I want to raise: "
        "should LEN plus CAB/RPV be the default for patients at high LTFU risk? We've seen at "
        "least two other patients started on CAB/RPV alone at other counties who showed up at "
        "MXM months later with new INSTI resistance. That's devastating. LEN might have prevented "
        "it. The downside -- LEN is a painful subcutaneous injection, harder to administer than "
        "IM CAB/RPV.")

    # ======================================================================
    # SLIDE 13: Case Divider - Patient 2
    # ======================================================================
    slide13 = add_divider_slide(prs, "Patient 2",
                                 "Update case from Nov 2024 Grand Rounds. Street/sheltered homeless, meth use, PTSD.")
    set_speaker_notes(slide13,
        "Patient 2. Some of you may remember this case from Grand Rounds last year. This is an update.")

    # ======================================================================
    # SLIDE 14: Patient 2 Timeline
    # ======================================================================
    slide14 = add_navy_content_slide(prs, "Patient 2 \u2014 Clinical Timeline")
    add_timeline_figure(slide14)
    set_speaker_notes(slide14,
        "Patient 2 was started on Cabenuva at another clinic before transferring to MXM. At "
        "baseline he had NNRTI mutations -- K103N, V108I -- but no INSTI mutations. He initially "
        "suppressed, but at dose 6 his viral load bounced to 37,000. Before the genotype came "
        "back, we added LEN -- reasoning that he likely had resistance to one class but probably "
        "not both, so LEN would give us approximately two active agents. The genotype was worse "
        "than expected -- high-level resistance to both CAB and RPV. We stopped CAB/RPV and "
        "prescribed oral ART, but he told us flat out he cannot take pills. He disengaged. "
        "Effectively on LEN monotherapy. Then at week 38 he comes back for an abscess -- totally "
        "unrelated to HIV care -- and his viral load is undetectable. On what was functionally "
        "LEN monotherapy for months, including about 10 weeks where LEN was likely subtherapeutic. "
        "We restarted everything at week 42.")

    # ======================================================================
    # SLIDE 15: Patient 2 Update/Today
    # ======================================================================
    slide15 = add_white_content_slide(prs, "Patient 2 \u2014 Today (April 2026)")

    # Big stats
    add_stat_callout(slide15, 0.5, 1.1, 4.0, 1.5, "~195+", "total weeks on LEN\u00b1CAB/RPV",
                     stat_color=TEAL, label_color=DARK_TEXT)
    add_stat_callout(slide15, 5.0, 1.1, 4.0, 1.5, "88", "additional weeks\nviral suppression",
                     stat_color=GREEN, label_color=DARK_TEXT)

    teaching = [
        "This is 1 patient. He may be an outlier. NOT an endorsement of LEN monotherapy.",
        "LEN resistance mutations (e.g., M66I) confer ~1.5% of wild-type replication capacity",
        "CAPELLA trial: 14 LEN resistance cases at 104 wks; 7 re-suppressed after adherence counseling/OBR switch",
        "Connects to Patient 1: LEN's long tail protects INSTI class even in lapses",
    ]
    add_bullet_text(slide15, 0.75, 3.0, 11.8, 3.0, teaching,
                    font_size=SMALL_SIZE, font_color=DARK_TEXT, line_spacing=1.3)

    # Humility quote
    quote_box = slide15.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), Inches(5.5), Inches(10.0), Inches(0.8)
    )
    quote_box.fill.solid()
    quote_box.fill.fore_color.rgb = LIGHT_TEAL
    quote_box.line.color.rgb = TEAL
    quote_box.line.width = Pt(1)
    tf = quote_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = ("\"I don't fully understand why this regimen is holding. "
                "If anyone wants to discuss the pharmacology after, I'd genuinely welcome that.\"")
    run.font.name = FONT_BODY
    run.font.size = CAPTION_SIZE
    run.font.color.rgb = NAVY
    run.font.italic = True

    set_speaker_notes(slide15,
        "As of today, Patient 2 has been on this regimen for about 195 weeks total. Eighty-eight "
        "more weeks of suppression since the restart. No viremia. He's still getting CAB/RPV "
        "despite high-level resistance to both drugs -- we keep it as a hedge. I want to be very "
        "clear: this is one patient. He could be an outlier. This is not an endorsement of LEN "
        "monotherapy as a strategy. But the biology is interesting -- LEN resistance mutations "
        "appear to devastate viral fitness, down to about 1.5% of wild-type replication. I "
        "honestly don't fully understand why this is holding. If anyone in this room wants to "
        "talk pharmacology after, I'd genuinely welcome that.")

    # ======================================================================
    # SLIDE 16: Case Divider - Patient 3
    # ======================================================================
    slide16 = add_divider_slide(prs, "Patient 3",
                                 "NOT our patient. Street homeless, relocated from another county. Cabenuva q8-weeks.")
    set_speaker_notes(slide16,
        "Patient 3 is interesting because he's not our patient. He showed up.")

    # ======================================================================
    # SLIDE 17: Patient 3 Background
    # ======================================================================
    slide17 = add_white_content_slide(prs, "Patient 3 \u2014 Background")
    p3_bullets = [
        "Street homeless, recently relocated from another county (~4 hrs away)",
        "Methamphetamine use disorder, bipolar 2, history of DVTs/PE, SSTI",
        "Started CAB/RPV 600/900mg q8 weeks at out-of-county clinic 8 months ago",
        "ALL injections on time at home clinic until recently",
        "Admitted ZSFG for SSTI \u2014 hospital tried Cabenuva but patient AMA'd",
        "Labs during hospitalization (1 mo ago): VL undetectable, CD4 670",
        "Presents to MXM 2 weeks post-discharge \u2014 for wound care, NOT HIV care",
        "Last Cabenuva dose: 10 weeks ago (2 weeks overdue for q8-week schedule)",
        "Declines labs, declines oral ART, declines returning to home county",
        "Open to receiving LA-ART today \u2014 same-day only",
    ]
    add_bullet_text(slide17, 0.75, 1.3, 11.8, 5.5, p3_bullets,
                    font_size=CAPTION_SIZE, font_color=DARK_TEXT, line_spacing=1.15)
    set_speaker_notes(slide17,
        "Patient 3 is street homeless, recently relocated from another county about four hours "
        "away. He was doing well on q8-week Cabenuva at his home clinic -- all injections on "
        "time. He got admitted to ZSFG for a skin infection, the hospital tried to give him his "
        "Cabenuva dose but he left AMA before they could. He shows up at MXM two weeks later for "
        "wound care -- not for HIV. Via Care Everywhere we can see his last Cabenuva was 10 weeks "
        "ago. He's two weeks overdue. His last labs from the hospital a month ago showed undetectable "
        "viral load, CD4 670. He won't do labs today. Won't take oral ART. Won't go back to his "
        "home county. But he'll take an injection right now if we can do it.")

    # ======================================================================
    # SLIDE 18: Patient 3 Poll 3
    # ======================================================================
    slide18 = add_poll_slide(prs,
        "He will not accept oral ART or labs. What do you do?",
        [
            "Order and administer Cabenuva today without labs",
            "Nothing \u2014 not your patient; return to out-of-county PCP",
            "Help him get a phone, encourage return when less activated",
            "Give him an oral ART sample anyway",
            "Another reasonable approach? (open to audience)",
        ]
    )
    set_speaker_notes(slide18,
        "No single right answer. What would you do? He's only two weeks past his due date. "
        "His last VL was undetectable a month ago -- but a month ago, not two weeks ago, and "
        "that distinction matters. What we actually did on visit 1: C. We helped him get phone "
        "resources and made a deal -- come back, do labs, and we'll give you the injection same day.")

    # ======================================================================
    # SLIDE 19: Patient 3 Resolution
    # ======================================================================
    slide19 = add_white_content_slide(prs, "Patient 3 \u2014 Resolution")

    add_textbox(slide19, 0.75, 1.2, 6.0, 0.4, "Visit 2 (1 week later \u2014 now 22 days past due):",
                font_size=SUBHEADING_SIZE, font_color=TEAL, bold=True)
    v2_bullets = [
        "Returns, calmer but somewhat volatile",
        "Agrees to labs IF same-day injection guaranteed",
        "Now 22 days overdue (vs. 14 at visit 1)",
    ]
    add_bullet_text(slide19, 0.75, 1.7, 5.5, 1.5, v2_bullets,
                    font_size=CAPTION_SIZE, font_color=DARK_TEXT, line_spacing=1.2)

    add_textbox(slide19, 7.0, 1.2, 5.5, 0.4, "Decision:",
                font_size=SUBHEADING_SIZE, font_color=TEAL, bold=True)
    dec_bullets = [
        "Used 1 of 2 emergency Cabenuva doses (CAB/RPV 600/900mg)",
        "Labs drawn simultaneously",
        "Via 403B Emergency ART program \u2014 NOT standard supply",
        "Told: return in 1\u20132 days for VL; establish care at MXM",
    ]
    add_bullet_text(slide19, 7.0, 1.7, 5.5, 2.0, dec_bullets,
                    font_size=CAPTION_SIZE, font_color=DARK_TEXT, line_spacing=1.2)

    # Why this was hard
    add_textbox(slide19, 0.75, 3.8, 11.8, 0.4, "Why This Was Hard:",
                font_size=SUBHEADING_SIZE, font_color=NAVY, bold=True)
    hard_bullets = [
        "Don't want to set precedent that non-patients get same-day Cabenuva",
        "If he developed resistance during lapse, giving dose = functional monotherapy",
        "Counter: extremely high INSTI resistance risk if untreated; very unlikely he could get home county dose",
    ]
    add_bullet_text(slide19, 0.75, 4.3, 11.8, 2.0, hard_bullets,
                    font_size=CAPTION_SIZE, font_color=DARK_TEXT, line_spacing=1.2)

    # Closing line
    close_box = slide19.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2.0), Inches(5.8), Inches(9.0), Inches(0.6)
    )
    close_box.fill.solid()
    close_box.fill.fore_color.rgb = RGBColor(0xFF, 0xF8, 0xE1)
    close_box.line.color.rgb = GOLD
    close_box.line.width = Pt(1)
    tf = close_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "\"I still don't know if this was the right call. Ask me next year.\""
    run.font.name = FONT_BODY
    run.font.size = SMALL_SIZE
    run.font.color.rgb = NAVY
    run.font.italic = True

    set_speaker_notes(slide19,
        "He came back a week later. Now 22 days overdue -- worse than before. He agreed to "
        "labs if we'd guarantee the injection same day. We used one of our two emergency "
        "Cabenuva doses -- obtained through the 403B Emergency ART program, not regular "
        "supply. I still don't know if this was the right call. We didn't want to set a "
        "precedent. If he'd developed resistance during the lapse, that injection is "
        "functional monotherapy. But the alternative was sending him away with nothing, "
        "and he was at extremely high risk of INSTI resistance if we didn't treat. Ask me "
        "next year.")

    # ======================================================================
    # SLIDE 20: Case Divider - Patient 4
    # ======================================================================
    slide20 = add_divider_slide(prs, "Patient 4",
                                 "52M, street homeless, CD4 <35, AIDS. Open question \u2014 genuine consult.")
    set_speaker_notes(slide20,
        "Last case. This one is an open question. I'm genuinely asking for your input.")

    # ======================================================================
    # SLIDE 21a: Patient 4 Background (Demographics/OI History)
    # ======================================================================
    slide21a = add_white_content_slide(prs, "Patient 4 \u2014 Background")

    p4_left = [
        "52-year-old man, street homeless ~6 months, no family",
        "HIV/AIDS: CD4 <35 (2%), VL 255,000 c/mL",
        "Not on ART \u22655 years despite prior Triumeq",
        "\"To tell the truth, I threw the medicine away\"",
        "Recently expressing interest in LA-ART",
        "Schizophrenia (no meds), stimulant use, AUD + withdrawal seizures",
        "M184V on recent genotype; CrAg NEGATIVE; afebrile",
    ]
    add_bullet_text(slide21a, 0.75, 1.3, 5.8, 4.0, p4_left,
                    font_size=CAPTION_SIZE, font_color=DARK_TEXT, line_spacing=1.15)

    add_textbox(slide21a, 7.0, 1.2, 5.8, 0.4, "Prior Opportunistic Infections:",
                font_size=SUBHEADING_SIZE, font_color=RED_ACCENT, bold=True)
    oi_bullets = [
        "PCP pneumonia (Dec 2025)",
        "Strep pneumo bacteremia + H. influenzae empyema (Mar 2025)",
        "\u2192 septic shock, bilateral chest tubes, pressors",
        "Esophageal candidiasis",
        "CMV viremia",
    ]
    add_bullet_text(slide21a, 7.0, 1.7, 5.8, 3.0, oi_bullets,
                    font_size=CAPTION_SIZE, font_color=DARK_TEXT, line_spacing=1.15)

    set_speaker_notes(slide21a,
        "Patient 4. Fifty-two years old, street homeless, CD4 less than 35. He's had PCP, "
        "Strep bacteremia with empyema, septic shock, cardiac tamponade from a pericardial "
        "effusion -- and a cardiac arrest. He threw his oral meds away. But recently he's "
        "been interested in LA-ART.")

    # ======================================================================
    # SLIDE 21b: Patient 4 Pericardial History + IRIS Question
    # ======================================================================
    slide21b = add_white_content_slide(prs, "Patient 4 \u2014 Pericardial History & IRIS Risk")

    add_textbox(slide21b, 0.75, 1.2, 5.8, 0.4, "Pericardial History (CRITICAL):",
                font_size=SUBHEADING_SIZE, font_color=RED_ACCENT, bold=True)
    peri_bullets = [
        "Mar 2025: Strep bacteremia \u2192 pericardial effusion \u2192 tamponade \u2192 PEA arrest \u2192 emergent drain",
        "Etiology: presumed infectious, never definitively established, NOT known TB",
        "Apr 2025: effusive-constrictive pericarditis, loculated effusion, pericardial mass",
        "CT surgery: too high risk; patient/family declined",
        "May 2025: AMA discharge",
        "Recent CT: NO current pericardial effusion (reassuring)",
    ]
    add_bullet_text(slide21b, 0.75, 1.7, 5.8, 3.5, peri_bullets,
                    font_size=Pt(13), font_color=DARK_TEXT, line_spacing=1.15)

    add_textbox(slide21b, 7.0, 1.2, 5.8, 0.4, "IRIS Risk Assessment:",
                font_size=SUBHEADING_SIZE, font_color=TEAL, bold=True)
    iris_bullets = [
        "Pericardial effusion IS a recognized IRIS manifestation (esp TB-IRIS)",
        "Cardiac tamponade as IRIS: documented but \"extremely rare\"",
        "CrAg negativity removes most dangerous trigger",
        "ART era: pericardial effusion rate 11% \u2192 0.25% (ART is protective)",
        "Concern is clinically legitimate but NOT a clear contraindication",
    ]
    add_bullet_text(slide21b, 7.0, 1.7, 5.8, 3.5, iris_bullets,
                    font_size=Pt(13), font_color=DARK_TEXT, line_spacing=1.15)

    # Bottom question
    q_box = slide21b.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.75), Inches(5.5), Inches(11.5), Inches(0.8)
    )
    q_box.fill.solid()
    q_box.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xCD)
    q_box.line.color.rgb = GOLD
    q_box.line.width = Pt(2)
    tf = q_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Clinical question: Risk of IRIS \u2192 new pericardial effusion \u2192 tamponade?"
    run.font.name = FONT_BODY
    run.font.size = Pt(18)
    run.font.color.rgb = NAVY
    run.font.bold = True

    set_speaker_notes(slide21b,
        "The clinical dilemma: his pericardial disease was incompletely treated. It was "
        "presumed infectious but never definitively characterized, and it's not TB. If we "
        "start ART and his immune system reconstitutes, could IRIS trigger a new pericardial "
        "effusion and tamponade? The literature says pericardial IRIS is recognized but "
        "extremely rare. His CrAg is negative, which is reassuring. And his recent CT shows "
        "no current effusion. But at a CD4 of 35, the biggest risk might be NOT treating.")

    # ======================================================================
    # SLIDE 22: Patient 4 Poll 4
    # ======================================================================
    slide22 = add_poll_slide(prs,
        "You are consulted on starting ART. What do you do?",
        [
            "Start Cabenuva (CAB/RPV) now",
            "Start Cabenuva + Lenacapavir now",
            "Start DRV/c/TAF/FTC (oral) now",
            "Start BIC/TAF/FTC (oral) now",
            "Wait \u2014 get LA-antipsychotic first, then reassess ART",
        ]
    )
    # Add discussion annotation below options
    add_textbox(slide22, 1.2, 6.2, 11.0, 0.5,
                "No correct answer \u2014 genuine consult. Current lean: E (LA-antipsychotic first), "
                "but CD4 <35 makes delay dangerous.",
                font_size=Pt(13), font_color=GOLD, italic=True)
    set_speaker_notes(slide22,
        "Last poll. What would you do? I'm leaning toward E right now -- if we could get him "
        "on a long-acting antipsychotic like paliperidone palmitate, his engagement and "
        "decision-making might improve enough to safely initiate ART with close monitoring. "
        "But I fully acknowledge the danger of delay at a CD4 of 35. Oral options are "
        "pharmacologically reasonable but this is a man who tells you to your face that he "
        "throws the medicine away. If we go LA-ART, the argument for adding LEN is strong "
        "-- high LTFU risk, M184V on genotype. I genuinely don't know the right answer here. "
        "I'd welcome your thoughts.")

    # ======================================================================
    # SLIDE 23: Key Takeaways / Thank You (split into 24 total with Patient 4 split)
    # ======================================================================
    slide23 = add_navy_content_slide(prs, "Key Takeaways")

    takeaways = [
        "LA-ART works in this population \u2014 100% suppression \u2014 but requires intensive infrastructure. Not set-it-and-forget-it.",
        "For patients at high LTFU risk, LEN + CAB/RPV may offer population-level harm reduction \u2014 protecting the INSTI class even when CAB levels drop.",
        "In low-barrier settings: act when the patient is in front of you. Waiting for perfect lab data can mean waiting forever.",
        "Starting ART in patients with AIDS and incompletely treated infectious pericarditis raises novel IRIS concerns the field has not yet addressed.",
        "Humility. We are working in genuinely new territory. These cases don't have clean endings \u2014 and that's the point.",
    ]

    for i, takeaway in enumerate(takeaways):
        y = 1.3 + i * 1.05
        # Number circle
        num_circle = slide23.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.75), Inches(y), Inches(0.4), Inches(0.4)
        )
        num_circle.fill.solid()
        num_circle.fill.fore_color.rgb = TEAL
        num_circle.line.fill.background()
        tf = num_circle.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(i + 1)
        run.font.name = FONT_BODY
        run.font.size = Pt(14)
        run.font.color.rgb = WHITE
        run.font.bold = True

        # Takeaway text
        add_textbox(slide23, 1.3, y - 0.05, 11.2, 0.5, takeaway,
                    font_name=FONT_BODY, font_size=CAPTION_SIZE,
                    font_color=WHITE)

    # Citations at bottom
    citations = (
        "Mehtani NJ et al. Open Forum Infect Dis. 2025;12(6):ofaf330.  |  "
        "Mehtani NJ et al. JAIDS. 2024;96:61-7.  |  "
        "LA-PrEP abstract, submitted 2026.  |  "
        "Oral vs. LA-ART comparison, in prep."
    )
    add_textbox(slide23, 0.75, 6.5, 11.8, 0.5, citations,
                font_size=Pt(10), font_color=MID_GRAY, italic=True,
                alignment=PP_ALIGN.CENTER)

    # Thank you
    add_textbox(slide23, 4.0, 7.0, 5.0, 0.4, "Thank you.",
                font_name=FONT_HEADING, font_size=Pt(24),
                font_color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)

    set_speaker_notes(slide23,
        "Five takeaways. One -- LA-ART works, 100% suppression, but it takes a team, outreach "
        "nurses, street medicine, showing up at SROs and shelters. Two -- lenacapavir may be "
        "the most important harm reduction tool we have for protecting the INSTI class at a "
        "population level. Three -- when someone is standing in front of you, that IS the window. "
        "Don't close it waiting for labs. Four -- we're encountering genuinely novel clinical "
        "territory with IRIS and pericardial disease that nobody has written about. And five -- "
        "humility. These cases don't have clean endings. That's the point. Thank you.")

    # ======================================================================
    # Save
    # ======================================================================
    prs.save(output_path)
    slide_count = len(prs.slides)
    print(f"Generated {output_path} with {slide_count} slides")
    return slide_count


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Generate HIV Grand Rounds PPTX")
    parser.add_argument("--output", default="hiv_grand_rounds_mxm.pptx",
                        help="Output filename (default: hiv_grand_rounds_mxm.pptx)")
    parser.add_argument("--template", default=None,
                        help="Path to UCSF template PPTX")
    args = parser.parse_args()

    # Find template
    script_dir = os.path.dirname(os.path.abspath(__file__))
    repo_root = os.path.dirname(os.path.dirname(script_dir))  # talks/49_*/.. -> talks/.. -> repo root
    template_path = args.template or os.path.join(
        repo_root, "examples", "test-files", "UCSF_ZSFG_Template_16x9.pptx"
    )

    if not os.path.exists(template_path):
        print(f"ERROR: Template not found at {template_path}", file=sys.stderr)
        sys.exit(1)

    output_path = os.path.join(script_dir, args.output)
    slide_count = build_deck(template_path, output_path)
    print(f"File size: {os.path.getsize(output_path) / 1024:.0f} KB")
