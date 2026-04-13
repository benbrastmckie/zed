#!/usr/bin/env python3
"""
Build HIV Grand Rounds PPTX Presentation
==========================================
Long-Acting ART & PrEP at MXM -- Nicky Mehtani, MD MPH
UCSF/ZSFG HIV Grand Rounds, April 2026

Generates a 23-slide PPTX with UCSF institutional theming.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# Theme constants
# ---------------------------------------------------------------------------
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
TEAL = RGBColor(0x00, 0x95, 0xA8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
ORANGE = RGBColor(0xE6, 0x6A, 0x00)
PURPLE = RGBColor(0x6A, 0x1B, 0x9A)
GOLD = RGBColor(0xC8, 0x96, 0x16)
GOLD_BG = RGBColor(0xFF, 0xF3, 0xCD)

HEADING_FONT = "Garamond"
BODY_FONT = "Arial"

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# ---------------------------------------------------------------------------
# Helper utilities
# ---------------------------------------------------------------------------

def _set_slide_bg(slide, color):
    """Set solid background color on a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, left, top, width, height, text, font_name=BODY_FONT,
                 font_size=Pt(22), font_color=WHITE, bold=False, italic=False,
                 alignment=PP_ALIGN.LEFT, word_wrap=True):
    """Add a text box with styled text. Returns the shape."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = font_size
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.italic = italic
    p.alignment = alignment
    return txBox


def _add_bullet_list(slide, left, top, width, height, items, font_name=BODY_FONT,
                     font_size=Pt(22), font_color=WHITE, bold_items=None,
                     spacing=Pt(6)):
    """Add a bulleted list. bold_items is a set of indices to bold."""
    if bold_items is None:
        bold_items = set()
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.name = font_name
        p.font.size = font_size
        p.font.color.rgb = font_color
        p.font.bold = i in bold_items
        p.space_after = spacing
        p.level = 0
    return txBox


def _add_stat_callout(slide, left, top, width, height, number_text, label_text,
                      number_color=TEAL, label_color=WHITE, number_size=Pt(54),
                      label_size=Pt(18), bg_color=None):
    """Add a large stat number with a label below it."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    # Number
    p = tf.paragraphs[0]
    p.text = number_text
    p.font.name = HEADING_FONT
    p.font.size = number_size
    p.font.color.rgb = number_color
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Label
    p2 = tf.add_paragraph()
    p2.text = label_text
    p2.font.name = BODY_FONT
    p2.font.size = label_size
    p2.font.color.rgb = label_color
    p2.font.bold = False
    p2.alignment = PP_ALIGN.CENTER

    return txBox


def _add_speaker_notes(slide, notes_text):
    """Add speaker notes to a slide."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes_text


def _add_callout_box(slide, left, top, width, height, text, bg_color=GOLD_BG,
                     text_color=DARK_GRAY, font_size=Pt(18), bold=False):
    """Add a rounded rectangle callout box with text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = BODY_FONT
    p.font.size = font_size
    p.font.color.rgb = text_color
    p.font.bold = bold
    p.alignment = PP_ALIGN.CENTER
    return shape


# ---------------------------------------------------------------------------
# Layout helpers (slide-level)
# ---------------------------------------------------------------------------

def add_teal_divider(prs, title, subtitle="", speaker_notes=""):
    """Teal background divider slide with centered title."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    _set_slide_bg(slide, TEAL)

    _add_textbox(slide, Inches(1), Inches(2.2), Inches(11.3), Inches(1.5),
                 title, font_name=HEADING_FONT, font_size=Pt(44),
                 font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    if subtitle:
        _add_textbox(slide, Inches(1), Inches(3.8), Inches(11.3), Inches(1),
                     subtitle, font_name=BODY_FONT, font_size=Pt(24),
                     font_color=WHITE, bold=False, alignment=PP_ALIGN.CENTER)

    _add_speaker_notes(slide, speaker_notes)
    return slide


def add_navy_slide(prs, title, speaker_notes=""):
    """Navy background slide with title. Returns slide for custom content."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, NAVY)

    _add_textbox(slide, Inches(0.7), Inches(0.4), Inches(11.9), Inches(0.9),
                 title, font_name=HEADING_FONT, font_size=Pt(32),
                 font_color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)

    _add_speaker_notes(slide, speaker_notes)
    return slide


def add_navy_bullet(prs, title, bullets, speaker_notes="", bold_items=None):
    """Navy background slide with title and bullet list."""
    slide = add_navy_slide(prs, title, speaker_notes)
    _add_bullet_list(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5),
                     bullets, font_color=WHITE, bold_items=bold_items)
    return slide


def add_navy_two_column(prs, title, left_title, left_items, right_title,
                        right_items, speaker_notes=""):
    """Navy two-column slide."""
    slide = add_navy_slide(prs, title, speaker_notes)

    # Left column header
    _add_textbox(slide, Inches(0.7), Inches(1.4), Inches(5.5), Inches(0.6),
                 left_title, font_name=BODY_FONT, font_size=Pt(22),
                 font_color=TEAL, bold=True)
    _add_bullet_list(slide, Inches(0.7), Inches(2.0), Inches(5.5), Inches(5.0),
                     left_items, font_size=Pt(18), font_color=WHITE)

    # Right column header
    _add_textbox(slide, Inches(6.8), Inches(1.4), Inches(5.8), Inches(0.6),
                 right_title, font_name=BODY_FONT, font_size=Pt(22),
                 font_color=TEAL, bold=True)
    _add_bullet_list(slide, Inches(6.8), Inches(2.0), Inches(5.8), Inches(5.0),
                     right_items, font_size=Pt(18), font_color=WHITE)
    return slide


def add_white_slide(prs, title, speaker_notes=""):
    """White background slide with navy title. Returns slide for custom content."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)

    _add_textbox(slide, Inches(0.7), Inches(0.4), Inches(11.9), Inches(0.9),
                 title, font_name=HEADING_FONT, font_size=Pt(32),
                 font_color=NAVY, bold=True, alignment=PP_ALIGN.LEFT)

    _add_speaker_notes(slide, speaker_notes)
    return slide


def add_white_bullet(prs, title, bullets, speaker_notes="", bold_items=None):
    """White background slide with title and bullet list."""
    slide = add_white_slide(prs, title, speaker_notes)
    _add_bullet_list(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5),
                     bullets, font_color=NAVY, bold_items=bold_items)
    return slide


def add_white_two_column(prs, title, left_title, left_items, right_title,
                         right_items, speaker_notes=""):
    """White two-column slide."""
    slide = add_white_slide(prs, title, speaker_notes)

    _add_textbox(slide, Inches(0.7), Inches(1.4), Inches(5.5), Inches(0.6),
                 left_title, font_name=BODY_FONT, font_size=Pt(22),
                 font_color=TEAL, bold=True)
    _add_bullet_list(slide, Inches(0.7), Inches(2.0), Inches(5.5), Inches(5.0),
                     left_items, font_size=Pt(18), font_color=NAVY)

    _add_textbox(slide, Inches(6.8), Inches(1.4), Inches(5.8), Inches(0.6),
                 right_title, font_name=BODY_FONT, font_size=Pt(22),
                 font_color=TEAL, bold=True)
    _add_bullet_list(slide, Inches(6.8), Inches(2.0), Inches(5.8), Inches(5.0),
                     right_items, font_size=Pt(18), font_color=NAVY)
    return slide


def add_poll_slide(prs, question, options, speaker_notes=""):
    """Poll Everywhere placeholder slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)

    # Poll Everywhere header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = TEAL
    header.line.fill.background()

    _add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12.3), Inches(0.9),
                 "Poll Everywhere", font_name=BODY_FONT, font_size=Pt(28),
                 font_color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)

    # Question
    _add_textbox(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(1.2),
                 question, font_name=HEADING_FONT, font_size=Pt(28),
                 font_color=NAVY, bold=True, alignment=PP_ALIGN.LEFT)

    # Options
    y = Inches(3.0)
    for opt in options:
        _add_textbox(slide, Inches(1.2), y, Inches(10.5), Inches(0.55),
                     opt, font_name=BODY_FONT, font_size=Pt(20),
                     font_color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT)
        y += Inches(0.7)

    _add_speaker_notes(slide, speaker_notes)
    return slide


# ---------------------------------------------------------------------------
# Slide builders (Slides 1-23)
# ---------------------------------------------------------------------------

def build_slide_01_title(prs):
    """Slide 1: Title Divider."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, TEAL)

    _add_textbox(slide, Inches(1), Inches(1.8), Inches(11.3), Inches(1.5),
                 "Long-Acting ART & PrEP at MXM",
                 font_name=HEADING_FONT, font_size=Pt(44),
                 font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    _add_textbox(slide, Inches(1), Inches(3.5), Inches(11.3), Inches(0.6),
                 "Nicky Mehtani, MD MPH",
                 font_name=BODY_FONT, font_size=Pt(24),
                 font_color=WHITE, bold=False, alignment=PP_ALIGN.CENTER)

    _add_textbox(slide, Inches(1), Inches(4.3), Inches(11.3), Inches(0.6),
                 "HIV Grand Rounds  |  UCSF/ZSFG  |  April 2026",
                 font_name=BODY_FONT, font_size=Pt(20),
                 font_color=WHITE, bold=False, alignment=PP_ALIGN.CENTER)

    _add_speaker_notes(slide, (
        "Thanks Stephen. I'm going to walk through our experience with long-acting "
        "ART and PrEP at MXM -- the numbers, and then four patients who teach us "
        "something about how this works in practice. I have some polls for you, so "
        "get your phones ready."
    ))


def build_slide_02_la_art_demographics(prs):
    """Slide 2: LA-ART: Who We Serve (Two Column -- Navy)."""
    slide = add_navy_two_column(
        prs,
        title="LA-ART: Who We Serve",
        left_title="Demographics (n=34, Nov 2021 -- Apr 2026)",
        left_items=[
            "Median age: 40",
            "22 cis-men, 8 trans women/non-binary, 4 cis-women",
            "50% Black, 26% Hispanic, 24% non-Hispanic White",
            "10 street homeless, 16 sheltered/couch-surfing, 8 unstably housed",
        ],
        right_title="Clinical Complexity",
        right_items=[
            "24% opioid use disorder",
            "11 initiated with CD4 <200 (7 with CD4 <50)",
            "71% started with detectable viremia",
            "17/24 had NOT been suppressed in our system for 3+ years",
        ],
        speaker_notes=(
            "So this is who we're working with. Thirty-four people since November "
            "2021. The demographics probably look familiar to anyone who works in "
            "the Tenderloin or SoMa. What I want you to notice is the clinical "
            "complexity column -- eighty-five percent meth use disorder, a third "
            "with serious psychotic illness, and most of them had not been virally "
            "suppressed in our system for at least three years. These are not "
            "patients who failed because they didn't try. These are patients the "
            "system wasn't designed for."
        )
    )

    # Stat callouts on right side
    _add_stat_callout(slide, Inches(7.0), Inches(4.6), Inches(2.8), Inches(1.4),
                      "85%", "methamphetamine use disorder",
                      number_color=TEAL, label_color=WHITE, number_size=Pt(48))
    _add_stat_callout(slide, Inches(10.0), Inches(4.6), Inches(2.8), Inches(1.4),
                      "32%", "schizophrenia or bipolar\nw/ psychosis",
                      number_color=TEAL, label_color=WHITE, number_size=Pt(48))


def build_slide_03_la_art_outcomes(prs):
    """Slide 3: LA-ART: Delivery & Outcomes (Navy)."""
    slide = add_navy_slide(prs, "LA-ART: Delivery & Outcomes",
        speaker_notes=(
            "Seven hundred forty-seven injection visits. Eighty percent happen at "
            "clinic, but seventeen percent happen out in the community -- shelters, "
            "syringe access programs, mobile outreach. Three percent in jails or "
            "hospitals. Ninety-three percent on time. And then the number that "
            "matters most: one hundred percent virally suppressed. Every single "
            "person. That includes patients who were not suppressed for years before "
            "starting. We have one person currently lost to follow-up -- you'll hear "
            "about him in a minute as Patient 1."
        ))

    # Big 100% stat callout
    _add_stat_callout(slide, Inches(3.5), Inches(1.5), Inches(6.3), Inches(2.2),
                      "100%", "virally suppressed at latest follow-up",
                      number_color=TEAL, label_color=WHITE, number_size=Pt(72))

    # 93% secondary
    _add_stat_callout(slide, Inches(0.8), Inches(4.0), Inches(3.5), Inches(1.3),
                      "93%", "on time (within 7-day window)",
                      number_color=TEAL, label_color=WHITE, number_size=Pt(44))

    # Delivery bullets
    _add_bullet_list(slide, Inches(5.0), Inches(4.2), Inches(7.8), Inches(3.0), [
        "747 total injection visits",
        "80% at MXM clinic | 17% mobile outreach | 3% jails/hospitals",
        "Median time on LA-ART: 116 weeks (IQR 60-145)",
        "9 successfully transferred to new primary care homes",
        "Current: 12 CAB/RPV q1-mo | 11 q2-mo | 9 CAB+/-RPV+LEN | 2 CAB/RPV+LEN",
    ], font_size=Pt(16), font_color=WHITE)


def build_slide_04_oral_vs_la_art(prs):
    """Slide 4: Oral vs. LA-ART Comparison (White two-column with table)."""
    slide = add_white_slide(prs, 'Oral vs. LA-ART Outcomes (2023-2024)',
        speaker_notes=(
            "This is from our manuscript in preparation. In 2023 and 2024, we had "
            "24 people on LA-ART and 123 on oral ART at MXM. One hundred percent "
            "of LA-ART patients were ever suppressed, versus seventy-two percent "
            "oral. At last follow-up, one hundred versus fifty-six. Now -- I need "
            "to be really clear about the caveat. These groups are not randomized. "
            "We selected patients for LA-ART who we could reliably find and who "
            "were interested. So the gap reflects both the intervention and the "
            "selection. But even with that caveat -- fifty-six percent versus one "
            "hundred percent is striking. And look at the encounter rate -- "
            "twenty-three visits per person-year. This is an intensive intervention. "
            "It is not set-it-and-forget-it."
        ))

    # Source line
    _add_textbox(slide, Inches(0.8), Inches(1.2), Inches(8.0), Inches(0.4),
                 "Mehtani et al., manuscript in preparation",
                 font_name=BODY_FONT, font_size=Pt(14),
                 font_color=DARK_GRAY, italic=True)

    # Build table
    rows, cols = 4, 3
    tbl_shape = slide.shapes.add_table(rows, cols, Inches(1.5), Inches(2.0),
                                        Inches(10.0), Inches(2.8))
    table = tbl_shape.table

    # Headers
    headers = ["Outcome", "Oral ART (n=123)", "LA-ART (n=24)"]
    data = [
        ["Ever virally suppressed", "72%", "100%"],
        ["Suppressed at last follow-up", "56%", "100%"],
        ["Mean encounters/person-year", "9.1", "23.0"],
    ]

    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.name = BODY_FONT
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY

    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.name = BODY_FONT
                p.font.size = Pt(20)
                p.font.color.rgb = NAVY
                p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            # Highlight 100% in teal
            if val == "100%":
                for p in cell.text_frame.paragraphs:
                    p.font.color.rgb = TEAL
                    p.font.bold = True
                    p.font.size = Pt(24)

    # Caveat
    _add_callout_box(slide, Inches(1.5), Inches(5.2), Inches(10.0), Inches(1.2),
                     ("CAVEAT: Groups not randomized. LA-ART patients selected for "
                      "reliable locatability and injection interest. Gap reflects "
                      "BOTH intervention AND selection."),
                     bg_color=RGBColor(0xFF, 0xF8, 0xE1), text_color=DARK_GRAY,
                     font_size=Pt(16), bold=False)


def build_slide_05_la_prep(prs):
    """Slide 5: LA-PrEP at MXM (White)."""
    slide = add_white_slide(prs, "LA-PrEP at MXM (Aug 2022 -- Dec 2025)",
        speaker_notes=(
            "Switching to PrEP. Sixty-eight people prescribed, fifty-two initiated. "
            "Similar population -- mostly unhoused, majority stimulant use disorder. "
            "Ninety percent retention at three months, sixty percent at six. Two "
            "seroconversions -- and this is important -- both happened during "
            "prolonged lapses in coverage, not during delayed dosing. So imperfect "
            "timing is not the same as failure. I also want to flag that we lost "
            "DPH-provided gift card incentives for about six to seven months starting "
            "last September, and that likely affected some of the retention data. "
            "The key message: continuity of coverage matters more than protocol "
            "perfection. For this population, we need to prioritize getting LA-PrEP "
            "into bodies -- rapid initiation, outreach-based delivery, flexible lab "
            "protocols."
        ))

    # Source line
    _add_textbox(slide, Inches(0.8), Inches(1.2), Inches(8.0), Inches(0.4),
                 "Mehtani et al., abstract submitted 2026",
                 font_name=BODY_FONT, font_size=Pt(14),
                 font_color=DARK_GRAY, italic=True)

    # Left side bullets
    _add_bullet_list(slide, Inches(0.8), Inches(1.8), Inches(7.0), Inches(4.5), [
        "Cohort: n=68 prescribed, 52 initiated (76%)",
        "65% unhoused | 55% stimulant use | 27% OUD | 19% trans/NB",
        "Delivery: 304 injections (58% clinic, 21% shelter, 23% outreach)",
        "Timing: 78% on time | 17% delayed 8-28d | 7% delayed >28d",
        "HIV testing: 68% included any testing; 61% guideline-concordant",
        "Seroconversions: 2 (both during prolonged lapses, NOT delayed dosing)",
        "Key message: Continuity of coverage > protocol perfection",
    ], font_size=Pt(18), font_color=NAVY, bold_items={6})

    # Stat callouts right
    _add_stat_callout(slide, Inches(8.5), Inches(2.0), Inches(4.0), Inches(1.5),
                      "90%", "retained at 3 months",
                      number_color=TEAL, label_color=NAVY, number_size=Pt(54))
    _add_stat_callout(slide, Inches(8.5), Inches(3.8), Inches(4.0), Inches(1.5),
                      "2", "seroconversions\n(prolonged lapses only)",
                      number_color=TEAL, label_color=NAVY, number_size=Pt(54))


def build_slide_06_patient1_divider(prs):
    """Slide 6: Patient 1 Case Divider."""
    add_teal_divider(prs, "Patient 1",
                     "Lenacapavir as Resistance Protection",
        speaker_notes=(
            "Let's get into cases. Patient 1 -- this is someone we published on "
            "last year, as Patient B in our Open Forum ID case series. He teaches "
            "us something about why lenacapavir might matter for population-level "
            "harm reduction."
        ))


def build_slide_07_patient1_background(prs):
    """Slide 7: Patient 1 Background (Navy)."""
    slide = add_navy_two_column(
        prs,
        title="Patient 1 -- Background",
        left_title="Demographics",
        left_items=[
            "49-year-old man, housed in SRO (Tenderloin)",
            "HIV diagnosed 2010",
            "Schizophrenia, bipolar disorder, psychosis NOS",
            "Methamphetamine use disorder",
            "Recurrent skin infections",
            "NOT interested in antipsychotics or psychiatry",
            "No cell phone -- agreed to outreach RN visits at SRO",
        ],
        right_title="HIV History & Baseline",
        right_items=[
            "Unable to maintain oral ART for 5 years",
            "Prior: EFV/TDF/FTC, RPV/TDF/FTC, EVG/c/TAF/FTC, BIC/TAF/FTC",
            "Hep B: immune",
            "CD4: 334 | VL: 340 c/mL",
            "3rd visit to MXM in 3 months -- first sustained engagement",
        ],
        speaker_notes=(
            "Forty-nine-year-old man, lives in an SRO in the Tenderloin. HIV since "
            "2010. Schizophrenia, bipolar, psychosis -- he is not interested in "
            "antipsychotics or psychiatry, and that is his right. Meth use disorder. "
            "No cell phone, but he agreed to let our outreach RN come to his SRO. "
            "He had not been able to stay on oral ART for five years -- tried four "
            "different regimens. When he came to MXM, his viral load was only 340 "
            "-- he'd been intermittently taking something -- and his CD4 was 334. "
            "But the key here is the genotype: Y181C plus M184V. That means "
            "high-level rilpivirine resistance. No INSTI mutations. So cabotegravir "
            "is still active, but rilpivirine is not."
        )
    )

    # Resistance callout box
    _add_callout_box(slide, Inches(6.8), Inches(5.3), Inches(5.8), Inches(1.0),
                     "Genotype: Y181C + M184V = RPV excluded. CAB remains active.",
                     bg_color=GOLD_BG, text_color=DARK_GRAY,
                     font_size=Pt(18), bold=True)


def build_slide_08_patient1_poll1(prs):
    """Slide 8: Patient 1 Poll 1."""
    add_poll_slide(prs,
        question="Would you start this patient on LA-ART?",
        options=[
            "A: No -- not with untreated schizophrenia/psychosis",
            "B: Yes, but wait until he has a cell phone",
            "C: Yes -- start CAB/RPV monthly",
            "D: Yes -- start CAB monthly + LEN q6 months",
            "E: Yes -- start CAB/RPV q2-months + LEN q6 months",
        ],
        speaker_notes=(
            "Take a minute. Get your phones out. Would you start this patient on "
            "long-acting ART, and if so, what regimen? [Pause for poll responses]"
        ))


def build_slide_09_patient1_answer(prs):
    """Slide 9: Patient 1 Poll 1 Answer + Clinical Course (Navy)."""
    slide = add_navy_slide(prs, "Patient 1 -- Answer & Clinical Course",
        speaker_notes=(
            "The answer is D. Drop the rilpivirine -- it's not doing anything with "
            "Y181C resistance. Start CAB monthly plus lenacapavir every six months. "
            "The lenacapavir acts as a capsid inhibitor backbone. If he lapses on "
            "CAB and CAB levels drop, lenacapavir is still there protecting the "
            "INSTI class. And for fourteen months, it worked beautifully. Sixteen "
            "injections, all on time, undetectable throughout. Then he disappeared. "
            "Three months later he shows up -- says he biked more than three hundred "
            "miles to another city and just got back. He's ten weeks late for CAB "
            "and four weeks late for LEN. He wants to restart."
        ))

    # Answer highlight
    _add_callout_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.7),
                     "Answer: D -- CAB monthly + LEN q6 months",
                     bg_color=RGBColor(0x1B, 0x5E, 0x20), text_color=WHITE,
                     font_size=Pt(22), bold=True)

    # Rationale bullets
    _add_bullet_list(slide, Inches(0.8), Inches(2.4), Inches(5.5), Inches(3.5), [
        "A: Untreated psych is the rule at MXM, not the exception",
        "B: No phone does not equal barrier -- outreach RN goes to SRO",
        "C: Y181C = RPV resistance. CAB alone risks INSTI resistance",
        "D: Drop RPV, add LEN as capsid inhibitor backbone",
        "E: Including RPV with resistance weakens regimen",
    ], font_size=Pt(16), font_color=WHITE, bold_items={3})

    # Clinical course on right
    _add_textbox(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(0.5),
                 "Clinical Course", font_name=BODY_FONT, font_size=Pt(22),
                 font_color=TEAL, bold=True)
    _add_bullet_list(slide, Inches(6.8), Inches(2.2), Inches(5.8), Inches(4.5), [
        "Phase 1 -- Success (14 months):",
        "  CAB monthly + LEN q6 months",
        "  16 injections, all on time, VL undetectable",
        "",
        "Then: LTFU",
        "  Disappeared for 3 months",
        "  Returns 10 wks late for CAB, 4 wks late for LEN",
        "  Reports: biked >300 miles to another city",
        "  Wants to restart injections",
    ], font_size=Pt(16), font_color=WHITE)


def build_slide_10_patient1_poll2(prs):
    """Slide 10: Patient 1 Poll 2."""
    add_poll_slide(prs,
        question=("In addition to drawing HIV VL + RNA genotype, "
                  "what do you do today?"),
        options=[
            "A: Nothing today -- have him return in 2-3 days once labs result",
            "B: Start BIC/TAF/FTC oral bridge until labs return",
            "C: Start DRV/c/TAF/FTC empirically until labs return",
            "D: Administer LEN + CAB today and plan for close follow-up",
            "E: Another reasonable option? (open to audience)",
        ],
        speaker_notes=(
            "Here's poll number two. He's standing in front of you right now. "
            "You've drawn labs -- VL and genotype. What do you do TODAY? This one "
            "doesn't have a clean answer. [Pause for poll responses]"
        ))


def build_slide_11_patient1_resolution(prs):
    """Slide 11: Patient 1 Resolution + LEN Pharmacology (Navy)."""
    slide = add_navy_slide(prs, "Patient 1 -- Resolution",
        speaker_notes=(
            "What we did: we gave him LEN and CAB that day. The window of "
            "engagement IS the moment of opportunity. If we send him away for "
            "labs, he might bike three hundred miles again tomorrow. Labs came "
            "back -- undetectable. Despite being ten weeks late for CAB. Two "
            "teaching points. First: in low-barrier settings, you act when the "
            "patient is in front of you. Waiting for perfect data means waiting "
            "forever. Second: lenacapavir likely protected him. LEN has a longer "
            "PK tail than CAB. When CAB went subtherapeutic, LEN was still "
            "working -- different mechanism, capsid inhibitor. Without LEN, we "
            "would almost certainly have seen INSTI resistance during that "
            "ten-week gap. He restarted, did another sixteen months perfectly. "
            "We published this case last year."
        ))

    # VL undetectable stat
    _add_stat_callout(slide, Inches(4.0), Inches(1.3), Inches(5.5), Inches(1.5),
                      "VL UNDETECTABLE", "despite 10 wks late CAB, 4 wks late LEN",
                      number_color=TEAL, label_color=WHITE, number_size=Pt(36))

    # Teaching point boxes
    _add_callout_box(slide, Inches(0.8), Inches(3.2), Inches(5.5), Inches(2.0),
                     ("Teaching Point 1: Low-Barrier Philosophy\n\n"
                      "In low-barrier settings, act with the opportunity "
                      "in front of you. Don't wait for perfect lab data -- "
                      "that delay = months of lost contact."),
                     bg_color=RGBColor(0x1A, 0x3A, 0x5C), text_color=WHITE,
                     font_size=Pt(16))

    _add_callout_box(slide, Inches(6.8), Inches(3.2), Inches(5.8), Inches(2.0),
                     ("Teaching Point 2: Why LEN Protected Him\n\n"
                      "LEN has a longer PK tail than CAB. When CAB went "
                      "subtherapeutic, LEN (capsid inhibitor) continued to "
                      "suppress. Without LEN, INSTI resistance was almost "
                      "certain during a 10-week CAB gap."),
                     bg_color=RGBColor(0x1A, 0x3A, 0x5C), text_color=WHITE,
                     font_size=Pt(16))

    # Phase 2 note
    _add_textbox(slide, Inches(0.8), Inches(5.6), Inches(11.5), Inches(0.7),
                 ("Phase 2 (restart): 16 more injections over ~16 months, all on "
                  "time, VL undetectable  |  Published: Mehtani et al., Open Forum "
                  "Infect Dis 2025 (PMC12188974)"),
                 font_name=BODY_FONT, font_size=Pt(14),
                 font_color=RGBColor(0xBB, 0xBB, 0xBB))


def build_slide_12_patient1_2nd_ltfu(prs):
    """Slide 12: Patient 1 2nd LTFU + Population Question (Navy)."""
    slide = add_navy_slide(prs, "Patient 1 -- Current Status & Population Question",
        speaker_notes=(
            "And then he disappeared again. Eleven weeks out from his last CAB "
            "dose. We think he's more than three thousand miles away based on a "
            "Care Everywhere note from an ED visit. He's the only patient in our "
            "cohort who's currently lost to follow-up. But here's the silver "
            "lining -- if resistance develops during this lapse, it will likely "
            "be to lenacapavir, not to the INSTI class. When he re-engages "
            "somewhere, CAB-based options may still be available. And this raises "
            "a bigger question. We have seen at least two other patients started "
            "on CAB/RPV alone at other counties who became LTFU and came to MXM "
            "with new INSTI resistance. That is devastating. Should we be adding "
            "lenacapavir as default for anyone at high LTFU risk -- as a harm "
            "reduction strategy to protect the INSTI class at a population level?"
        ))

    # Current status bullets
    _add_bullet_list(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.5), [
        "2nd LTFU after 16 months of success",
        "11 weeks since last CAB dose",
        "Still within LEN window (due ~26 weeks)",
        "Cannot locate -- suspected >3,000 miles away",
        "Only patient in our cohort currently LTFU",
    ], font_size=Pt(18), font_color=WHITE)

    # Silver lining
    _add_bullet_list(slide, Inches(0.8), Inches(4.2), Inches(5.5), Inches(1.5), [
        "If resistance develops: likely to LEN, not INSTI class",
        "INSTI-based options may remain available on re-engagement",
    ], font_size=Pt(16), font_color=RGBColor(0xBB, 0xBB, 0xBB))

    # Population question in large teal text
    _add_textbox(slide, Inches(6.5), Inches(1.8), Inches(6.3), Inches(3.5),
                 ("Should LEN + CAB/RPV become the DEFAULT for high LTFU-risk "
                  "patients -- as harm reduction to protect the INSTI class?"),
                 font_name=HEADING_FONT, font_size=Pt(28),
                 font_color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)

    # Supporting data note
    _add_textbox(slide, Inches(6.5), Inches(5.5), Inches(6.3), Inches(1.2),
                 ("At least 2 other patients at MXM started on CAB/RPV alone at "
                  "other counties, became LTFU, presented months later with new "
                  "INSTI resistance. LEN might have prevented it."),
                 font_name=BODY_FONT, font_size=Pt(14),
                 font_color=RGBColor(0xBB, 0xBB, 0xBB))


def build_slide_13_patient2_divider(prs):
    """Slide 13: Patient 2 Case Divider."""
    add_teal_divider(prs, "Patient 2",
                     "Update: Unintentional LEN Monotherapy",
        speaker_notes=(
            "Patient 2. Some of you may remember this case from Grand Rounds "
            "last November. This is an update -- here's where things stand now."
        ))


def build_slide_14_patient2_timeline(prs):
    """Slide 14: Patient 2 Timeline (custom shapes on navy)."""
    slide = add_navy_slide(prs, "Patient 2 -- Treatment Timeline",
        speaker_notes=(
            "Patient 2 -- street homeless, meth use disorder, PTSD, psychosis. "
            "Started CAB/RPV at another clinic before MXM. Baseline VL almost "
            "178,000, K103N on genotype. Initially he suppresses, but then at "
            "dose six his viral load comes back at 37,000. We don't have the "
            "genotype yet, but we're pretty sure he's developed resistance to "
            "one class -- either CAB or RPV. So we add lenacapavir that same day, "
            "thinking we still have two active agents. Then the genotype comes "
            "back and it's devastating -- high-level resistance to both CAB and "
            "RPV. We stop them, prescribe oral ART. He tells us he cannot take "
            "pills. He never picks them up. He's on LEN monotherapy. He comes "
            "back ten months later for an abscess and his viral load is "
            "undetectable. We restart everything. As of today, he's been "
            "suppressed for about a hundred and ninety-five weeks total."
        ))

    # Timeline horizontal arrow
    arrow_y = Inches(3.5)
    arrow_left = Inches(0.5)
    arrow_right = Inches(12.8)

    # Draw arrow line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, arrow_left, arrow_y, arrow_right - arrow_left, Pt(4)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = WHITE
    line.line.fill.background()

    # Timeline points with colored circles
    timeline_points = [
        # (x_pos, color, label_above, label_below)
        (Inches(1.0), PURPLE, "W-29\nBaseline", "VL 177,992\nCD4 187"),
        (Inches(2.8), GREEN, "Dose 2-3", "VL 33\nthen <30"),
        (Inches(4.2), ORANGE, "Dose 4-5\n(delayed)", "12+21d late"),
        (Inches(5.8), GREEN, "W0: Dose 6\n+ ADD LEN", "VL 37,464"),
        (Inches(7.8), RGBColor(0xCC, 0x00, 0x00), "Geno returns\nSTOP CAB/RPV", "HIGH-LEVEL\nINSTI + NNRTI R"),
        (Inches(9.5), TEAL, "W38\nReturns", "VL <30"),
        (Inches(10.8), GREEN, "W42\nRestart all", "VL undetectable"),
        (Inches(12.0), GREEN, "TODAY\n~195+ wks", "Suppressed\nCD4 554"),
    ]

    circle_size = Inches(0.35)
    for x, color, label_above, label_below in timeline_points:
        # Circle
        circ = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x - circle_size // 2,
            arrow_y - circle_size // 2,
            circle_size, circle_size
        )
        circ.fill.solid()
        circ.fill.fore_color.rgb = color
        circ.line.fill.background()

        # Label above
        _add_textbox(slide, x - Inches(0.6), arrow_y - Inches(1.6),
                     Inches(1.4), Inches(1.1),
                     label_above, font_name=BODY_FONT, font_size=Pt(10),
                     font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

        # Label below
        _add_textbox(slide, x - Inches(0.6), arrow_y + Inches(0.4),
                     Inches(1.4), Inches(1.0),
                     label_below, font_name=BODY_FONT, font_size=Pt(9),
                     font_color=RGBColor(0xBB, 0xBB, 0xBB),
                     alignment=PP_ALIGN.CENTER)

    # Gold callout for key decision
    _add_callout_box(slide, Inches(4.5), Inches(5.5), Inches(4.5), Inches(0.8),
                     "W0: Added LEN pre-genotype -- thought 2 active agents remained",
                     bg_color=GOLD_BG, text_color=DARK_GRAY,
                     font_size=Pt(13), bold=False)

    _add_callout_box(slide, Inches(9.2), Inches(5.5), Inches(3.8), Inches(0.8),
                     "Patient refused oral ART -- effectively LEN monotherapy",
                     bg_color=GOLD_BG, text_color=DARK_GRAY,
                     font_size=Pt(13), bold=False)


def build_slide_15_patient2_teaching(prs):
    """Slide 15: Patient 2 Teaching Points (Navy)."""
    slide = add_navy_slide(prs, "Patient 2 -- Update & Teaching Points",
        speaker_notes=(
            "I want to be really clear about this case. This is one patient. He "
            "might be an outlier. I am not endorsing lenacapavir monotherapy. But "
            "the possible mechanism is interesting -- LEN resistance mutations "
            "confer very low viral fitness. The virus might develop resistance but "
            "essentially can't replicate well enough to matter. I genuinely don't "
            "understand why this is working as well as it is. If anyone wants to "
            "talk pharmacology after, I would welcome that conversation."
        ))

    # Stat callout
    _add_stat_callout(slide, Inches(8.5), Inches(1.3), Inches(4.0), Inches(1.5),
                      "195+", "weeks on regimen\ncontinuously suppressed",
                      number_color=TEAL, label_color=WHITE, number_size=Pt(54))

    # Teaching points
    _add_bullet_list(slide, Inches(0.8), Inches(1.6), Inches(7.2), Inches(5.0), [
        "1. This is 1 patient. He may be an outlier. Do NOT interpret as "
        "endorsement of LEN monotherapy.",
        "",
        "2. Possible mechanism: LEN resistance mutations (e.g., M66I) confer very "
        "low viral fitness (~1.5% of wild-type replication). Resistance may evolve "
        "but with poor fitness.",
        "",
        "3. CAPELLA trial: 14 LEN resistance cases at 104 weeks; 7 re-suppressed "
        "after adherence counseling or OBR switch.",
        "",
        "4. Connects to Patient 1: LEN's long tail seems to protect INSTI class "
        "even in lapses.",
    ], font_size=Pt(16), font_color=WHITE)

    # Closing quote
    _add_textbox(slide, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.8),
                 ('"I don\'t fully understand why this regimen is holding. '
                  'If anyone wants to discuss the pharmacology after, I\'d '
                  'genuinely welcome that."'),
                 font_name=BODY_FONT, font_size=Pt(16),
                 font_color=RGBColor(0xBB, 0xBB, 0xBB), italic=True)


def build_slide_16_patient3_divider(prs):
    """Slide 16: Patient 3 Case Divider."""
    add_teal_divider(prs, "Patient 3",
                     "Not Our Patient",
        speaker_notes=(
            "Patient 3. This one is about a situation I think a lot of us are "
            "going to face more and more as long-acting ART expands beyond "
            "specialty clinics."
        ))


def build_slide_17_patient3_background(prs):
    """Slide 17: Patient 3 Background (Navy)."""
    slide = add_navy_two_column(
        prs,
        title="Patient 3 -- Background",
        left_title="Demographics & History",
        left_items=[
            "Street homeless in SF (recently from another county ~4 hrs away)",
            "Methamphetamine use disorder, bipolar 2",
            "History of DVTs/PE (not on anticoagulation), recurrent SSTI",
            "Started CAB/RPV q8-weeks at out-of-county clinic 8 months ago",
            "On time for ALL injections at home clinic until recently",
        ],
        right_title="Recent Events",
        right_items=[
            "Admitted ZSFG for SSTI, IV antibiotics",
            "Hospital tried to give Cabenuva -- couldn't confirm last date",
            "Patient left AMA before dose given",
            "Labs during admission: VL undetectable, CD4 670",
            "",
            "Presents to MXM ~2 wks after discharge for wound eval",
            "NOT for HIV care",
            "Last injection: 10 weeks ago (2 weeks overdue)",
            "Declines labs, oral ART, returning to home county",
            "Open to receiving LA-ART today -- same-day only",
        ],
        speaker_notes=(
            "Patient 3. Street homeless, recently came to SF from about four "
            "hours away. He was started on Cabenuva q-eight-weeks at his home "
            "county clinic and was on time for every injection until he moved. "
            "He got admitted to General for a skin infection, the hospital tried "
            "to give his Cabenuva but couldn't confirm the date of his last "
            "injection, and he left AMA. Two weeks later he shows up at MXM -- "
            "not for HIV care, just for wound eval. He's not our patient. His "
            "last injection was ten weeks ago, so he's two weeks overdue. His "
            "last labs from the hospital a month ago were undetectable. But he "
            "declines labs today. Declines oral ART. Declines going back to his "
            "home county. He says: just give me my injection."
        )
    )

    # Emphasis callout
    _add_callout_box(slide, Inches(3.5), Inches(6.0), Inches(6.3), Inches(0.7),
                     "NOT our patient  |  10 weeks since last dose  |  Declines everything except same-day injection",
                     bg_color=GOLD_BG, text_color=DARK_GRAY,
                     font_size=Pt(16), bold=True)


def build_slide_18_patient3_poll(prs):
    """Slide 18: Patient 3 Poll 3."""
    add_poll_slide(prs,
        question=("Recognizing he will not accept oral ART or labs -- "
                  "what do you do?"),
        options=[
            "A: Order and administer Cabenuva today without labs",
            "B: Nothing -- not your patient, send to home county PCP",
            "C: Help him get a phone and encourage return when less activated",
            "D: Give him an oral ART sample anyway",
            "E: Another reasonable approach? (open to audience)",
        ],
        speaker_notes=(
            "What do you do? He will not accept labs. He will not accept oral "
            "ART. He will not go back to his home county. What he will accept "
            "is an injection, right now. [Pause for poll responses]"
        ))


def build_slide_19_patient3_resolution(prs):
    """Slide 19: Patient 3 Resolution (Navy)."""
    slide = add_navy_slide(prs, "Patient 3 -- Resolution",
        speaker_notes=(
            "What we actually did on visit one was C -- we sent him away with "
            "resources to get a phone and a deal: come back, get labs and an "
            "injection the same day. He came back a week later, now twenty-two "
            "days overdue. We used one of our two emergency Cabenuva doses -- "
            "these are from the 403B Emergency ART program, not standard supply. "
            "We drew labs simultaneously and told him he needs to establish care "
            "formally. Was this the right call? I honestly don't know. We didn't "
            "want to set a precedent that non-patients can show up for same-day "
            "Cabenuva. But he was at real risk of INSTI resistance if we did "
            "nothing. Ask me next year."
        ))

    # Visit 1
    _add_textbox(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.5),
                 "Visit 1", font_name=BODY_FONT, font_size=Pt(22),
                 font_color=TEAL, bold=True)

    _add_callout_box(slide, Inches(0.8), Inches(1.9), Inches(5.5), Inches(0.6),
                     "Answer: C -- sent him away with phone resources and a deal",
                     bg_color=RGBColor(0x1B, 0x5E, 0x20), text_color=WHITE,
                     font_size=Pt(18), bold=True)

    _add_bullet_list(slide, Inches(0.8), Inches(2.7), Inches(5.5), Inches(2.5), [
        "A: Only 2 weeks past due, case for acting is real",
        "B: No car, no money, no intention to return = abandonment",
        "C: Buys time without abandoning or making a risky call",
        "D: He explicitly declined; pressing risks losing him",
    ], font_size=Pt(15), font_color=WHITE)

    # Visit 2
    _add_textbox(slide, Inches(6.8), Inches(1.4), Inches(5.8), Inches(0.5),
                 "Visit 2 (1 week later -- now 22 days past due)",
                 font_name=BODY_FONT, font_size=Pt(22),
                 font_color=TEAL, bold=True)

    _add_bullet_list(slide, Inches(6.8), Inches(2.1), Inches(5.8), Inches(3.0), [
        "Returns calmer, still no home county access",
        "Agrees to labs IF same-day injection guaranteed",
        "Used 1 of 2 emergency Cabenuva doses (403B program)",
        "Labs drawn simultaneously",
        "Told him: return in 1-2 days, must establish formal care",
    ], font_size=Pt(16), font_color=WHITE)

    # Closing quote
    _add_textbox(slide, Inches(2.5), Inches(5.8), Inches(8.5), Inches(1.0),
                 '"I still don\'t know if this was the right call. Ask me next year."',
                 font_name=HEADING_FONT, font_size=Pt(26),
                 font_color=TEAL, italic=True, alignment=PP_ALIGN.CENTER)


def build_slide_20_patient4_divider(prs):
    """Slide 20: Patient 4 Case Divider."""
    add_teal_divider(prs, "Patient 4",
                     "AIDS, Pericarditis, and IRIS",
        speaker_notes=(
            "Last case. This one is genuinely an open question. I'm asking for "
            "your help."
        ))


def build_slide_21_patient4_background(prs):
    """Slide 21: Patient 4 Background (Navy, dense -- may benefit from 2 slides)."""
    slide = add_navy_two_column(
        prs,
        title="Patient 4 -- Background",
        left_title="Demographics & HIV/AIDS",
        left_items=[
            "52-year-old man, street homeless in SF ~6 months",
            "Not on ART for 5+ years despite prior Triumeq",
            '"To tell the truth, I threw the medicine away"',
            "Recently expressing interest in LA-ART",
            "M184V on recent genotype",
            "Schizophrenia (not on meds), stimulant use, AUD + withdrawal seizures",
            "Recent PE (2024, not on anticoagulation)",
        ],
        right_title="Prior OIs & Pericardial History",
        right_items=[
            "PCP pneumonia (Dec 2025)",
            "Strep pneumo bacteremia + H. flu empyema (Mar 2025)",
            "  -- Septic shock, bilateral chest tubes",
            "Esophageal/mucocutaneous candidiasis",
            "Mar 2025: Strep -> pericardial effusion -> tamponade -> PEA arrest",
            "  -- Emergent drain; culture neg, 2700 WBC, PMN predominant",
            "  -- Etiology: presumed infectious, never established",
            "  -- CT surgery: too high risk (AIDS, adherence)",
            "Recent CT chest: NO current pericardial effusion",
        ],
        speaker_notes=(
            "Fifty-two-year-old man. Street homeless, no family. CD4 less than "
            "thirty-five, viral load two hundred fifty-five thousand. He has not "
            "been on ART for at least five years. When we prescribe pills, he "
            "throws them away -- his words, not mine. But recently he's been "
            "saying he's interested in long-acting injections. His history is "
            "extraordinary. PCP pneumonia, strep bacteremia with empyema and "
            "septic shock requiring bilateral chest tubes, and then -- the "
            "critical piece -- pericardial effusion that went to tamponade and "
            "PEA arrest. He survived. The pericardial fluid was culture negative "
            "but looked infectious. CT surgery felt it was too risky to do a "
            "pericardial window given his AIDS and likely adherence issues. His "
            "recent CT shows no current effusion, which is reassuring. But the "
            "question you're being asked: should we start LA-ART?"
        )
    )

    # Key stat callouts
    _add_stat_callout(slide, Inches(0.8), Inches(6.0), Inches(3.0), Inches(1.0),
                      "CD4 <35", "(2%)",
                      number_color=RGBColor(0xFF, 0x60, 0x60), label_color=WHITE,
                      number_size=Pt(36), label_size=Pt(16))
    _add_stat_callout(slide, Inches(4.0), Inches(6.0), Inches(3.5), Inches(1.0),
                      "VL 255,000", "c/mL",
                      number_color=RGBColor(0xFF, 0x60, 0x60), label_color=WHITE,
                      number_size=Pt(36), label_size=Pt(16))

    # Question callout
    _add_callout_box(slide, Inches(7.8), Inches(6.0), Inches(5.2), Inches(1.0),
                     "IRIS risk if starting ART at CD4 <35 with incompletely treated infectious pericarditis?",
                     bg_color=GOLD_BG, text_color=DARK_GRAY,
                     font_size=Pt(15), bold=True)


def build_slide_22_patient4_poll(prs):
    """Slide 22: Patient 4 Poll 4 + Discussion."""
    slide = add_poll_slide(prs,
        question="You are consulted on starting ART. What do you do?",
        options=[
            "A: Start Cabenuva (CAB/RPV) now",
            "B: Start Cabenuva + Lenacapavir now",
            "C: Start DRV/c/TAF/FTC (oral) now",
            "D: Start BIC/TAF/FTC (oral) now",
            "E: Wait -- pursue LA-antipsychotic first, then reassess ART",
        ],
        speaker_notes=(
            "Last poll. You're being consulted. What do you do? [Pause for "
            "responses] I'll tell you where I'm leaning -- option E. If we can "
            "get him on paliperidone palmitate, a long-acting antipsychotic, his "
            "engagement and decision-making might improve enough to make ART "
            "initiation safer and more sustainable. But I acknowledge the danger "
            "of waiting at a CD4 under thirty-five. Pericardial IRIS is a "
            "recognized phenomenon, especially with TB, though his CrAg is "
            "negative which removes the most dangerous trigger. And in the ART "
            "era, pericardial effusion rates dropped from eleven percent to a "
            "quarter of a percent -- ART is generally protective for pericardial "
            "disease. I genuinely don't know the right answer here. I'd welcome "
            "your thoughts."
        ))

    # Discussion points below poll options
    _add_textbox(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.6),
                 ("IRIS context: CrAg neg (reassuring) | Pericardial IRIS "
                  "documented but rare | ART-era pericardial effusion: 11% to "
                  "0.25% | CMV viremia: additional theoretical trigger"),
                 font_name=BODY_FONT, font_size=Pt(12),
                 font_color=DARK_GRAY, italic=True)


def build_slide_23_takeaways(prs):
    """Slide 23: Takeaways / Thank You."""
    slide = add_white_slide(prs, "Key Takeaways",
        speaker_notes=(
            "Five takeaways. One: this works. One hundred percent suppression. "
            "But it is intensive -- twenty-three encounters per person-year, "
            "outreach nurses going to SROs and shelters. Two: lenacapavir may "
            "be doing something important at the population level, protecting "
            "the INSTI class when patients lapse. Three: act when the patient "
            "is in front of you. Four: we're encountering genuinely new clinical "
            "questions -- like IRIS risk in patients with AIDS and incompletely "
            "treated pericarditis -- that the field hasn't addressed yet. And "
            "five: humility. These cases don't have clean endings, and that's "
            "the point. Thank you. Happy to take questions."
        ))

    takeaways = [
        ('1. LA-ART works in this population -- ',
         '100% suppression',
         ' -- but requires intensive infrastructure. Not set-it-and-forget-it.'),
        ('2. For patients at high LTFU risk, ',
         'LEN + CAB/RPV may offer population-level harm reduction',
         ' -- protecting the INSTI class even when CAB levels drop.'),
        ('3. In low-barrier settings: ',
         'act when the patient is in front of you.',
         ' Waiting for perfect lab data can mean waiting forever.'),
        ('4. Starting ART in patients with AIDS and incompletely treated '
         'infectious pericarditis raises ',
         'novel IRIS concerns',
         ' the field has not yet addressed.'),
        ('5. ',
         'Humility.',
         ' We are working in genuinely new territory. These cases don\'t '
         'have clean endings -- and that\'s the point.'),
    ]

    y = Inches(1.6)
    for prefix, bold_part, suffix in takeaways:
        txBox = slide.shapes.add_textbox(Inches(0.8), y, Inches(11.5), Inches(0.9))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        # Add runs with mixed formatting
        run1 = p.add_run()
        run1.text = prefix
        run1.font.name = BODY_FONT
        run1.font.size = Pt(18)
        run1.font.color.rgb = NAVY

        run2 = p.add_run()
        run2.text = bold_part
        run2.font.name = BODY_FONT
        run2.font.size = Pt(18)
        run2.font.color.rgb = TEAL
        run2.font.bold = True

        run3 = p.add_run()
        run3.text = suffix
        run3.font.name = BODY_FONT
        run3.font.size = Pt(18)
        run3.font.color.rgb = NAVY

        y += Inches(1.0)

    # Citations
    citations = [
        "1. Mehtani NJ et al. Open Forum Infect Dis. 2025;12(6):ofaf330. PMC12188974.",
        "2. Mehtani NJ et al. J Acquir Immune Defic Syndr. 2024;96:61-7. PMC11009050.",
        "3. Mehtani NJ et al. LA-PrEP abstract, submitted 2026.",
        "4. Mehtani NJ et al. Oral vs. LA-ART manuscript, in preparation.",
    ]
    _add_bullet_list(slide, Inches(0.8), Inches(6.3), Inches(8.0), Inches(1.0),
                     citations, font_size=Pt(11), font_color=DARK_GRAY,
                     spacing=Pt(2))

    # Thank you
    _add_textbox(slide, Inches(9.0), Inches(6.3), Inches(4.0), Inches(0.8),
                 "Thank you.\nQuestions?",
                 font_name=HEADING_FONT, font_size=Pt(28),
                 font_color=TEAL, bold=True, alignment=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------------------
# Main build function
# ---------------------------------------------------------------------------

def build_presentation():
    """Build the complete 23-slide presentation."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Build all 23 slides in order
    build_slide_01_title(prs)           # Slide 1: Title
    build_slide_02_la_art_demographics(prs)  # Slide 2: LA-ART Demographics
    build_slide_03_la_art_outcomes(prs)      # Slide 3: LA-ART Outcomes
    build_slide_04_oral_vs_la_art(prs)       # Slide 4: Oral vs LA-ART
    build_slide_05_la_prep(prs)              # Slide 5: LA-PrEP
    build_slide_06_patient1_divider(prs)     # Slide 6: Patient 1 Divider
    build_slide_07_patient1_background(prs)  # Slide 7: Patient 1 Background
    build_slide_08_patient1_poll1(prs)       # Slide 8: Patient 1 Poll 1
    build_slide_09_patient1_answer(prs)      # Slide 9: Patient 1 Answer
    build_slide_10_patient1_poll2(prs)       # Slide 10: Patient 1 Poll 2
    build_slide_11_patient1_resolution(prs)  # Slide 11: Patient 1 Resolution
    build_slide_12_patient1_2nd_ltfu(prs)    # Slide 12: Patient 1 2nd LTFU
    build_slide_13_patient2_divider(prs)     # Slide 13: Patient 2 Divider
    build_slide_14_patient2_timeline(prs)    # Slide 14: Patient 2 Timeline
    build_slide_15_patient2_teaching(prs)    # Slide 15: Patient 2 Teaching
    build_slide_16_patient3_divider(prs)     # Slide 16: Patient 3 Divider
    build_slide_17_patient3_background(prs)  # Slide 17: Patient 3 Background
    build_slide_18_patient3_poll(prs)        # Slide 18: Patient 3 Poll
    build_slide_19_patient3_resolution(prs)  # Slide 19: Patient 3 Resolution
    build_slide_20_patient4_divider(prs)     # Slide 20: Patient 4 Divider
    build_slide_21_patient4_background(prs)  # Slide 21: Patient 4 Background
    build_slide_22_patient4_poll(prs)        # Slide 22: Patient 4 Poll
    build_slide_23_takeaways(prs)            # Slide 23: Takeaways

    return prs


def main():
    """Build and save the presentation."""
    script_dir = os.path.dirname(os.path.abspath(__file__))
    output_path = os.path.join(script_dir, "HIV_Grand_Rounds_MXM.pptx")

    print("Building HIV Grand Rounds presentation...")
    prs = build_presentation()

    prs.save(output_path)
    slide_count = len(prs.slides)
    file_size = os.path.getsize(output_path)

    print(f"Saved: {output_path}")
    print(f"Slides: {slide_count}")
    print(f"Size: {file_size:,} bytes ({file_size / 1024:.1f} KB)")

    # Verify speaker notes
    notes_count = 0
    for slide in prs.slides:
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            notes_count += 1
    print(f"Slides with speaker notes: {notes_count}/{slide_count}")


if __name__ == "__main__":
    main()
